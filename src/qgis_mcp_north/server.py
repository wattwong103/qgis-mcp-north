#!/usr/bin/env python3
"""qgis-mcp-north — focused QGIS MCP server for transportation research figures.

Forked from nkarasiak/qgis-mcp v0.2.1. Cut hard to 12 workflow tools + 1 escape
hatch (`qgis_eval`). See ``docs/DESIGN.md`` for tool-surface design, response
shapes, error taxonomy, and roadmap.

This module is the **v0.2 scaffold**: every tool is registered with FastMCP and
its full input/output schema, but every body raises ``NotImplementedError``.
Implementation lands tool-by-tool from v0.3 onward (see DESIGN.md §7).

Default socket port: 9877 (vs upstream nkarasiak/qgis-mcp on 9876). Both servers
can run side-by-side; the LLM picks per request based on tool descriptions.
"""

from __future__ import annotations

import logging
import os
import sys
from logging.handlers import RotatingFileHandler
from typing import Annotated, Literal

from mcp.server.fastmcp import FastMCP
from mcp.types import ToolAnnotations
from pydantic import BaseModel, Field

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------


def _setup_logging() -> logging.Logger:
    """stderr (WARNING+) plus optional rotating file handler."""
    log = logging.getLogger("QgisMcpNorthServer")
    log.handlers.clear()

    fmt = logging.Formatter("%(asctime)s - %(name)s - %(levelname)s - %(message)s")

    stderr_handler = logging.StreamHandler(sys.stderr)
    stderr_handler.setLevel(logging.WARNING)
    stderr_handler.setFormatter(fmt)
    log.addHandler(stderr_handler)

    default_log_file = os.path.join("~", ".local", "share", "qgis-mcp-north", "server.log")
    log_file_raw = os.environ.get("QGIS_MCP_NORTH_LOG_FILE", default_log_file)
    log_level_name = os.environ.get("QGIS_MCP_NORTH_LOG_LEVEL", "INFO").upper()
    file_level = getattr(logging, log_level_name, logging.INFO)

    if log_file_raw:
        log_file = os.path.expanduser(log_file_raw)
        os.makedirs(os.path.dirname(log_file), exist_ok=True)
        file_handler = RotatingFileHandler(log_file, maxBytes=5 * 1024 * 1024, backupCount=3)
        file_handler.setLevel(file_level)
        file_handler.setFormatter(fmt)
        log.addHandler(file_handler)
        log.setLevel(min(logging.WARNING, file_level))
    else:
        log.setLevel(logging.WARNING)
    return log


logger = _setup_logging()


# ---------------------------------------------------------------------------
# FastMCP server
# ---------------------------------------------------------------------------

SERVER_INSTRUCTIONS = """\
qgis-mcp-north — opinionated QGIS MCP for transportation research figure
pipelines (PFLOW, GUFM). Reach for these tools when you need to:
- Inspect or load a vector/raster layer.
- Render a choropleth from a zone polygon + value CSV.
- Render trajectories from CSV (PFLOW schema: lon, lat, datetime, trip_id).
- Render origin-destination flow arcs from an OD CSV + zones layer.
- Drop a batch of figures into a PowerPoint deck for a weekly report.
- Run arbitrary PyQGIS as an escape hatch (`qgis_eval`).

If you need general QGIS-API features (feature editing, processing
algorithms, layer-tree management, plugin tooling), use the upstream
nkarasiak/qgis-mcp server side-by-side. Both run together; pick per task.
"""

mcp = FastMCP("qgis-mcp-north", instructions=SERVER_INSTRUCTIONS)


# ---------------------------------------------------------------------------
# Pydantic response models — one per tool; see DESIGN.md §4 for field rationale
# ---------------------------------------------------------------------------


class FieldInfo(BaseModel):
    name: str
    type: str
    n_unique: int | None = None


class LayerInfo(BaseModel):
    """Read-only metadata about a vector or raster file on disk."""

    path: str
    geometry_type: Literal["point", "line", "polygon", "raster", "no_geom"]
    crs: str
    n_features: int
    extent: list[float] = Field(..., description="[xmin, ymin, xmax, ymax]")
    fields: list[FieldInfo]


class LoadedLayer(LayerInfo):
    """Layer that has been added to the active QGIS project."""

    layer_id: str


class LayerSummary(BaseModel):
    """Compact summary of a layer registered in a project."""

    layer_id: str
    name: str
    geometry_type: str
    visible: bool


class LayoutSummary(BaseModel):
    name: str


class ProjectInfo(BaseModel):
    project_path: str
    crs: str
    extent: list[float]
    layers: list[LayerSummary]
    layouts: list[LayoutSummary]


class ClassEntry(BaseModel):
    value: str
    color: str
    n_features: int


class StyleResult(BaseModel):
    layer_id: str
    n_classes: int
    classes: list[ClassEntry]


class GraduatedStyleResult(StyleResult):
    breaks: list[float]
    mode: str


class RenderResult(BaseModel):
    output_path: str
    width: int
    height: int
    dpi: int
    extent: list[float]
    crs: str
    n_layers: int


class JoinResult(BaseModel):
    csv: str
    field: str
    n_matched: int
    n_unmatched: int


class ChoroplethResult(RenderResult):
    field: str
    n_classes: int
    breaks: list[float]
    mode: str
    min_value: float
    max_value: float
    n_features: int
    join: JoinResult | None = None


class TrajectoryResult(RenderResult):
    n_trajectories: int
    n_points_total: int
    n_points_rendered: int
    downsampled: bool
    time_range: list[str] | None = None
    modes: list[str] | None = None
    used_movingpandas: bool


class ODFlowResult(RenderResult):
    n_flows: int
    n_flows_rendered: int
    n_zones: int
    max_flow: float
    min_flow_rendered: float
    n_unmatched_origins: int
    n_unmatched_destinations: int


class ExportResult(BaseModel):
    output_path: str
    format: str
    n_pages: int
    layout_name: str


class BatchManifestEntry(BaseModel):
    value: str
    output_path: str
    extent: list[float]


class BatchError(BaseModel):
    value: str
    error: str


class BatchRenderResult(BaseModel):
    output_dir: str
    n_rendered: int
    manifest: list[BatchManifestEntry]
    errors: list[BatchError]


class PptxResult(BaseModel):
    pptx_path: str
    n_slides_added: int
    n_slides_total: int
    slide_titles: list[str | None]


class EvalResult(BaseModel):
    stdout: str
    stderr: str
    return_values: dict | None = None
    exception: str | None = None


# ---------------------------------------------------------------------------
# Stub helper
# ---------------------------------------------------------------------------


def _stub(tool_name: str, design_section: str) -> None:
    """Raise a clear NotImplementedError pointing at the design doc."""
    raise NotImplementedError(
        f"{tool_name} is a v0.2 scaffold stub. Implementation lands in a future "
        f"version — see docs/DESIGN.md §{design_section} for the spec."
    )


_RASTER_EXTENSIONS = (".tif", ".tiff", ".geotiff", ".vrt", ".asc", ".img", ".jp2")


def _is_raster_path(path: str) -> bool:
    return path.lower().endswith(_RASTER_EXTENSIONS)


def _translate_geometry_type(plugin_type: str) -> str:
    """Plugin's ``vector_{0,1,2}`` / ``raster`` → DESIGN.md geometry_type enum."""
    if plugin_type == "raster":
        return "raster"
    if plugin_type.startswith("vector_"):
        idx = plugin_type.split("_", 1)[1]
        return {"0": "point", "1": "line", "2": "polygon"}.get(idx, "no_geom")
    return "no_geom"


def _load_and_get_info(executor, abs_path: str, name: str | None = None):
    """Load a layer + fetch metadata. Removes the layer if get_layer_info fails.

    Returns ``(layer_id, info_dict, is_raster)``. Layer stays loaded on success
    so the caller can decide whether to remove it (transient inspect) or keep it
    (persistent load).
    """
    is_raster = _is_raster_path(abs_path)
    load_cmd = "add_raster_layer" if is_raster else "add_vector_layer"
    params = {"path": abs_path}
    if name:
        params["name"] = name
    load_result = executor.dispatch(load_cmd, params)
    layer_id = load_result["id"]
    try:
        info = executor.dispatch("get_layer_info", {"layer_id": layer_id})
    except Exception:
        try:
            executor.dispatch("remove_layer", {"layer_id": layer_id})
        except Exception:
            logger.warning("post-error cleanup failed for %s", layer_id, exc_info=True)
        raise
    return layer_id, info, is_raster


def _layer_info_kwargs(abs_path: str, info: dict, is_raster: bool) -> dict:
    """Translate plugin's get_layer_info response into LayerInfo constructor kwargs."""
    extent_dict = info["extent"]
    fields = [
        FieldInfo(name=f["name"], type=f["type"], n_unique=None)
        for f in info.get("fields", [])
    ]
    return {
        "path": abs_path,
        "geometry_type": _translate_geometry_type(info["type"]),
        "crs": info["crs"],
        "n_features": 0 if is_raster else info.get("feature_count", 0),
        "extent": [
            extent_dict["xmin"], extent_dict["ymin"],
            extent_dict["xmax"], extent_dict["ymax"],
        ],
        "fields": fields,
    }


# ---------------------------------------------------------------------------
# Tools — Inspection & loading (3)
# ---------------------------------------------------------------------------


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=True, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_layer_inspect(
    path: Annotated[str, Field(description="Absolute path to a shapefile (.shp), GeoPackage (.gpkg), GeoJSON (.geojson), or raster (.tif).")],
) -> LayerInfo:
    """Read metadata for a layer file *without* loading it.

    When to use: before any render_* call, to confirm field names, geometry
    type, CRS, and feature count. Cheap (no project mutation, no QGIS render).

    Inputs: absolute path to a vector or raster file on disk.

    Returns: ``LayerInfo`` with path, geometry_type ∈ {point, line, polygon,
    raster, no_geom}, CRS string (EPSG:#### or proj string), n_features,
    extent as ``[xmin, ymin, xmax, ymax]``, and the list of fields with their
    types and unique-value counts.

    Chains into: ``qgis_load_layer`` (if you need to render),
    ``qgis_render_choropleth`` (pass ``zones_path=path``),
    ``qgis_render_trajectory`` (pass ``input_path=path``).
    """
    from qgis_mcp_north.executors import get_executor

    abs_path = os.path.abspath(path)
    executor = get_executor()
    layer_id, info, is_raster = _load_and_get_info(executor, abs_path)
    try:
        return LayerInfo(**_layer_info_kwargs(abs_path, info, is_raster))
    finally:
        try:
            executor.dispatch("remove_layer", {"layer_id": layer_id})
        except Exception:
            logger.warning("transient cleanup failed for %s", layer_id, exc_info=True)


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=False, destructiveHint=False, openWorldHint=True
    )
)
def qgis_load_layer(
    path: Annotated[str, Field(description="Absolute path to a vector or raster file.")],
    name: Annotated[str | None, Field(description="Display name for the layer in the project. Defaults to the file's stem.")] = None,
    crs: Annotated[str | None, Field(description='Override CRS, e.g. "EPSG:4326". Use when the file is missing a .prj sidecar.')] = None,
) -> LoadedLayer:
    """Add a layer to the active project and return its metadata + layer_id.

    When to use: when downstream tools need to operate on a registered layer
    (styling, rendering with multiple layers, batch operations). For one-shot
    figure renders on a single file, ``qgis_render_choropleth`` and friends
    accept a path directly and don't require this step.

    Returns: ``LoadedLayer`` — same fields as ``LayerInfo`` plus ``layer_id``
    used by ``qgis_style_*`` and ``qgis_render_map``.

    Chains into: ``qgis_style_categorized``, ``qgis_style_graduated``,
    ``qgis_render_map``.
    """
    if crs is not None:
        raise NotImplementedError(
            "crs override is a v0.4 feature — for v0.3, ensure the .prj is correct, "
            "or use qgis_eval to call layer.setCrs() after loading. "
            "Tracking: docs/DESIGN.md §8."
        )
    from qgis_mcp_north.executors import get_executor

    abs_path = os.path.abspath(path)
    executor = get_executor()
    layer_id, info, is_raster = _load_and_get_info(executor, abs_path, name=name)
    return LoadedLayer(layer_id=layer_id, **_layer_info_kwargs(abs_path, info, is_raster))


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=False, destructiveHint=False, openWorldHint=True
    )
)
def qgis_project_load(
    qgz_path: Annotated[str, Field(description="Absolute path to a .qgz or .qgs project file.")],
) -> ProjectInfo:
    """Load a saved QGIS project (``.qgz`` / ``.qgs``).

    When to use: when the user has already designed a map in QGIS — layers
    loaded, styles applied, layouts configured — and you want to export
    figures from it. The W17 weekly-deck pattern uses this: load project,
    then ``qgis_export_layout`` or ``qgis_render_map``.

    Returns: ``ProjectInfo`` listing all layers (with layer_ids) and all
    print-composer layouts available for export.

    Chains into: ``qgis_export_layout``, ``qgis_batch_render``,
    ``qgis_render_map``.
    """
    _stub("qgis_project_load", "4 / Inspection & Loading")


# ---------------------------------------------------------------------------
# Tools — Styling (2)
# ---------------------------------------------------------------------------


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=False
    )
)
def qgis_style_categorized(
    layer_id: Annotated[str, Field(description="layer_id from qgis_load_layer or qgis_project_load.")],
    field: Annotated[str, Field(description="Field name to categorize on (string-typed).")],
    palette: Annotated[str, Field(description='ColorBrewer palette name, e.g. "Set2", "Paired", "Dark2".')] = "Set2",
    classes: Annotated[list[str] | None, Field(description="Optional subset/order of category values to render; others get a default 'no data' style.")] = None,
) -> StyleResult:
    """Apply categorical (one-color-per-value) symbology to a vector layer.

    When to use: when ``field`` holds a small set of categorical values
    (e.g., ``transport_mode``, ``taxi_type``, ``zone_type``) and you want
    each category in a distinct color.

    Returns: ``StyleResult`` with the resolved class list and per-class
    feature counts.
    """
    _stub("qgis_style_categorized", "4 / Styling")


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=False
    )
)
def qgis_style_graduated(
    layer_id: Annotated[str, Field(description="layer_id from qgis_load_layer or qgis_project_load.")],
    field: Annotated[str, Field(description="Numeric field to bin on (e.g., total_trips, trip_count, vkt).")],
    n_classes: Annotated[int, Field(description="Number of bins.", ge=2, le=15)] = 5,
    mode: Annotated[Literal["quantile", "equal_interval", "natural_breaks", "pretty"], Field(description="Binning strategy.")] = "quantile",
    palette: Annotated[str, Field(description='Sequential colorbrewer palette, e.g. "YlOrRd", "Blues", "Viridis".')] = "YlOrRd",
) -> GraduatedStyleResult:
    """Apply graduated (value-based color ramp) symbology — the choropleth primitive.

    When to use: as a low-level building block. For zone-level choropleths,
    prefer ``qgis_render_choropleth`` (one call instead of three).

    Returns: ``GraduatedStyleResult`` — class list + per-class feature counts +
    the resolved ``breaks`` array + the ``mode`` used.
    """
    _stub("qgis_style_graduated", "4 / Styling")


# ---------------------------------------------------------------------------
# Tools — Rendering (4)
# ---------------------------------------------------------------------------


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_render_map(
    layer_ids: Annotated[list[str], Field(description="Layers to render, drawn in order bottom→top.")],
    output_png: Annotated[str, Field(description="Absolute path for the output PNG.")],
    width: Annotated[int, Field(description="Image width in pixels.", ge=200, le=8000)] = 1600,
    height: Annotated[int, Field(description="Image height in pixels.", ge=200, le=8000)] = 1200,
    dpi: Annotated[int, Field(description="Image DPI (affects font size).", ge=72, le=600)] = 150,
    extent: Annotated[list[float] | None, Field(description="Render extent as [xmin, ymin, xmax, ymax]. If omitted, uses the union of layer extents with 5% padding.")] = None,
    background: Annotated[str, Field(description='Map background, named or hex (e.g. "white", "#fafafa", "transparent").')] = "white",
) -> RenderResult:
    """Render a list of already-loaded layers to PNG.

    When to use: generic render. For domain-specific cases prefer the
    workflow tools (``qgis_render_choropleth``, ``qgis_render_trajectory``,
    ``qgis_render_od_flows``) — they handle data loading, styling, and
    extent inference for you.

    Returns: ``RenderResult`` with the absolute output_path, image dims,
    final extent, CRS, and number of layers rendered.

    Chains into: ``qgis_figures_to_pptx``, ``qgis_batch_render``.
    """
    from qgis_mcp_north.executors import get_executor

    abs_output = os.path.abspath(output_png)
    params: dict = {
        "layer_ids": list(layer_ids),
        "output_png": abs_output,
        "width": width,
        "height": height,
        "dpi": dpi,
        "background": background,
    }
    if extent is not None:
        params["extent"] = list(extent)

    result = get_executor().dispatch("render_layers_to_path", params, timeout=60)
    return RenderResult(
        output_path=result["output_path"],
        width=result["width"],
        height=result["height"],
        dpi=result["dpi"],
        extent=result["extent"],
        crs=result["crs"],
        n_layers=result["n_layers"],
    )


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_render_choropleth(
    zones_path: Annotated[str, Field(description="Absolute path to a polygon layer (shp, gpkg, geojson). For PFLOW prefecture choropleth, use polbnda_jpn_new.shp.")],
    value_field: Annotated[str, Field(description="Field name to render. If value_csv is given, this is the column in that CSV; else this is an attribute on zones_path.")],
    output_png: Annotated[str, Field(description="Absolute path for the output PNG.")],
    value_csv: Annotated[str | None, Field(description="Optional CSV path joined to zones_path on join_field. Use this for the PFLOW pattern: zone_trips.csv joined to a zones polygon.")] = None,
    join_field: Annotated[str, Field(description='Common key between zones_path and value_csv. PFLOW uses "zone_id" (MFS-coded), "PRF_CODE" (prefecture), or similar.')] = "zone_id",
    n_classes: Annotated[int, Field(description="Number of choropleth bins.", ge=2, le=15)] = 5,
    mode: Annotated[Literal["quantile", "equal_interval", "natural_breaks", "pretty"], Field(description="Binning strategy.")] = "quantile",
    palette: Annotated[str, Field(description='Sequential colorbrewer palette, e.g. "YlOrRd", "Blues", "Viridis".')] = "YlOrRd",
    title: Annotated[str | None, Field(description="Optional title rendered at the top of the figure.")] = None,
    legend: Annotated[bool, Field(description="Render a legend with class breaks.")] = True,
    basemap_paths: Annotated[list[str] | None, Field(description="Optional basemap layers drawn under the choropleth (e.g., coastline, rivers, prefecture borders).")] = None,
    width: Annotated[int, Field(description="Image width in pixels.", ge=200, le=8000)] = 1600,
    height: Annotated[int, Field(description="Image height in pixels.", ge=200, le=8000)] = 1200,
    dpi: Annotated[int, Field(description="Image DPI.", ge=72, le=600)] = 150,
) -> ChoroplethResult:
    """Render a zone-level choropleth in one call. PFLOW workflow tool.

    When to use: any zone-aggregated visualization where each polygon gets
    a color from a numeric value. Replaces ``qgis_load_layer`` +
    ``qgis_style_graduated`` + ``qgis_render_map`` chain.

    Two data shapes supported:
    1. ``value_field`` is already an attribute on ``zones_path`` → render directly.
    2. ``value_csv`` is provided → left-join ``value_csv[join_field]`` to
       ``zones_path[join_field]``, then render. Mismatches surface in the
       response as ``join.n_unmatched``, never as silent zero values.

    PFLOW example: ``zones_path=polbnda_jpn_new.shp``,
    ``value_csv=zone_trips.csv``, ``value_field=total_trips``,
    ``join_field=zone_id``.

    Returns: ``ChoroplethResult`` with output path, breaks, min/max, and a
    ``join`` block (if value_csv was used) reporting matched vs unmatched.

    Chains into: ``qgis_figures_to_pptx``, ``qgis_batch_render``.
    """
    import csv as _csv

    from qgis_mcp_north.errors import ExecutorError, FieldNotFoundError, JoinError
    from qgis_mcp_north.executors import get_executor

    abs_zones = os.path.abspath(zones_path)
    abs_output = os.path.abspath(output_png)
    abs_basemaps = [os.path.abspath(p) for p in (basemap_paths or [])]
    abs_csv = os.path.abspath(value_csv) if value_csv else None

    value_dict: dict[str, float] | None = None
    if abs_csv is not None:
        with open(abs_csv, encoding="utf-8", newline="") as f:
            reader = _csv.DictReader(f)
            csv_columns = reader.fieldnames or []
            if value_field not in csv_columns:
                raise FieldNotFoundError(value_field, csv_columns)
            if join_field not in csv_columns:
                raise FieldNotFoundError(join_field, csv_columns)
            value_dict = {}
            for row in reader:
                key = row[join_field]
                raw = row[value_field]
                try:
                    value_dict[str(key)] = float(raw)
                except (TypeError, ValueError):
                    logger.warning(
                        "skipping non-numeric value in %s: %s=%r → %s=%r",
                        abs_csv, join_field, key, value_field, raw,
                    )

    params: dict = {
        "zones_path": abs_zones,
        "value_field": value_field,
        "output_png": abs_output,
        "value_dict": value_dict,
        "join_field": join_field,
        "n_classes": n_classes,
        "mode": mode,
        "palette": palette,
        "basemap_paths": abs_basemaps,
        "width": width,
        "height": height,
        "dpi": dpi,
    }

    try:
        result = get_executor().dispatch("render_choropleth", params, timeout=60)
    except ExecutorError as err:
        if "JOIN_NO_MATCH" in str(err):
            raise JoinError(err.message) from err
        raise

    join_block = None
    if value_dict is not None:
        join_block = JoinResult(
            csv=abs_csv or "",
            field=join_field,
            n_matched=result["n_matched"],
            n_unmatched=result["n_unmatched"],
        )
    return ChoroplethResult(
        output_path=result["output_path"],
        width=result["width"],
        height=result["height"],
        dpi=result["dpi"],
        extent=result["extent"],
        crs=result["crs"],
        n_layers=result["n_layers"],
        field=result["field"],
        n_classes=result["n_classes"],
        breaks=result["breaks"],
        mode=result["mode"],
        min_value=result["min_value"],
        max_value=result["max_value"],
        n_features=result["n_features"],
        join=join_block,
    )


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_render_trajectory(
    input_path: Annotated[str, Field(description="Absolute path to a CSV (with lon, lat, datetime, trip_id columns by default — matches PFLOW trajectory schema) or GPX file.")],
    output_png: Annotated[str, Field(description="Absolute path for the output PNG.")],
    lon_col: Annotated[str, Field(description="Longitude column name in the CSV.")] = "lon",
    lat_col: Annotated[str, Field(description="Latitude column name in the CSV.")] = "lat",
    time_col: Annotated[str, Field(description="Datetime column name (ISO 8601 or PFLOW 'YYYY-MM-DD HH:MM:SS' string).")] = "datetime",
    id_col: Annotated[str, Field(description="Trajectory grouping column. Each unique value is one trajectory.")] = "trip_id",
    mode_col: Annotated[str | None, Field(description='Optional categorical column to color trajectories by (e.g., "transport_mode").')] = None,
    render_mode: Annotated[Literal["lines", "points", "heatmap"], Field(description="Visualization style.")] = "lines",
    sample_rate: Annotated[float, Field(description="Fraction of points to keep (1.0 = all, 0.01 = every 100th).", gt=0.0, le=1.0)] = 1.0,
    max_points: Annotated[int, Field(description="Hard cap on rendered points; exceeded → automatic downsample with response flag.", ge=1000)] = 500_000,
    basemap_paths: Annotated[list[str] | None, Field(description="Optional basemap layers drawn under trajectories.")] = None,
    extent: Annotated[list[float] | None, Field(description="[lon_min, lat_min, lon_max, lat_max] in EPSG:4326. Clips before rendering.")] = None,
    width: Annotated[int, Field(description="Image width in pixels.", ge=200, le=8000)] = 1600,
    height: Annotated[int, Field(description="Image height in pixels.", ge=200, le=8000)] = 1200,
    dpi: Annotated[int, Field(description="Image DPI.", ge=72, le=600)] = 150,
) -> TrajectoryResult:
    """Render trajectory data from CSV/GPX. PFLOW/GUFM workflow tool.

    When to use: visualizing GPS-style point sequences. Defaults match
    PFLOW trajectory CSV schema (``lon, lat, datetime, trip_id,
    transport_mode``). PFLOW files are large (3M+ rows, ~1 GB each); use
    ``sample_rate`` and ``extent`` to keep renders fast.

    If ``movingpandas`` and the QGIS Trajectools plugin are installed,
    auto-uses them for richer rendering (speed bins, stop detection); else
    falls back to plain line/point rendering. Reports which path was taken
    via ``used_movingpandas``.

    Returns: ``TrajectoryResult`` with totals, rendered count, downsample
    flag, observed time range, modes seen.

    Chains into: ``qgis_figures_to_pptx``.
    """
    _stub("qgis_render_trajectory", "4 / Rendering")


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_render_od_flows(
    od_csv: Annotated[str, Field(description="Absolute path to a long-format OD CSV. PFLOW schema: origin, destination, trip_count, avg_distance_km.")],
    zones_layer_path: Annotated[str, Field(description="Absolute path to a polygon layer with one feature per zone, keyed by zone_id_field.")],
    output_png: Annotated[str, Field(description="Absolute path for the output PNG.")],
    origin_col: Annotated[str, Field(description="Origin zone column in od_csv.")] = "origin",
    dest_col: Annotated[str, Field(description="Destination zone column in od_csv.")] = "destination",
    value_col: Annotated[str, Field(description="Flow magnitude column. Arc widths scale linearly with this.")] = "trip_count",
    zone_id_field: Annotated[str, Field(description="Zone identifier field on zones_layer_path. Must match origin_col / dest_col values.")] = "zone_id",
    top_n: Annotated[int | None, Field(description="Render only the top-N flows by value. None renders all matched flows.")] = None,
    basemap_paths: Annotated[list[str] | None, Field(description="Optional basemap layers drawn under arcs.")] = None,
    width: Annotated[int, Field(description="Image width in pixels.", ge=200, le=8000)] = 1600,
    height: Annotated[int, Field(description="Image height in pixels.", ge=200, le=8000)] = 1200,
    dpi: Annotated[int, Field(description="Image DPI.", ge=72, le=600)] = 150,
) -> ODFlowResult:
    """Render origin-destination arcs over a zones layer. PFLOW workflow tool.

    When to use: any "flows between regions" visualization. PFLOW example:
    ``od_csv=truck/run_*/od_flows.csv``,
    ``zones_layer_path=polbnda_jpn_new.shp``, ``zone_id_field=PRF_CODE``.

    PFLOW reality: multiple zone-id systems coexist (MFS##, PRF##, Z##) in
    different files. The caller must align ``od_csv[origin_col / dest_col]``
    with ``zones_layer[zone_id_field]``. Mismatches surface as
    ``n_unmatched_origins`` / ``n_unmatched_destinations`` — loud, not
    silent zero-flow renders.

    Returns: ``ODFlowResult`` with rendered count, max flow, and unmatched
    counts on each side.

    Chains into: ``qgis_figures_to_pptx``.
    """
    _stub("qgis_render_od_flows", "4 / Rendering")


# ---------------------------------------------------------------------------
# Tools — Export & Batch & Delivery (3)
# ---------------------------------------------------------------------------


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_export_layout(
    qgz_path: Annotated[str, Field(description="Absolute path to a .qgz or .qgs project containing the layout.")],
    layout_name: Annotated[str, Field(description="Print-composer layout name. List available layouts via qgis_project_load.")],
    output_path: Annotated[str, Field(description="Absolute path for the output file. Extension can be .png, .pdf, or .svg.")],
    format: Annotated[Literal["png", "pdf", "svg"], Field(description="Output format. Should match output_path extension.")] = "png",
    dpi: Annotated[int, Field(description="Export DPI.", ge=72, le=600)] = 300,
) -> ExportResult:
    """Export a saved print-composer layout to file.

    When to use: when the user has designed a polished figure layout in
    QGIS (titles, legends, scale bars, multi-panel) and wants to export it
    rather than re-render programmatically.

    Returns: ``ExportResult`` with format, page count, and resolved layout name.

    Chains into: ``qgis_figures_to_pptx``, ``qgis_batch_render``.
    """
    _stub("qgis_export_layout", "4 / Export & Batch & Delivery")


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_batch_render(
    template_qgz: Annotated[str, Field(description="Absolute path to a template project. Must contain a single 'active' layer to filter and (optionally) a layout to export.")],
    attribute: Annotated[str, Field(description="Field on the active layer to filter by.")],
    values: Annotated[list[str], Field(description="Filter values to iterate. One render per value.")],
    output_dir: Annotated[str, Field(description="Absolute path to a directory where renders are written.")],
    layout_name: Annotated[str | None, Field(description="If given, exports the layout with this name; otherwise renders the map canvas.")] = None,
    filename_template: Annotated[str, Field(description="Filename template, e.g. '{value}.png' or 'choropleth_{value}.png'.")] = "{value}.png",
) -> BatchRenderResult:
    """Fan-out: render the same template per filter value. Workflow tool.

    When to use: 'render the OD map for each scenario', 'render the
    choropleth for each timestep', 'one figure per prefecture'. Emits a
    manifest you can feed straight into ``qgis_figures_to_pptx``.

    Returns: ``BatchRenderResult`` with output_dir, n_rendered, manifest
    of (value → output_path → extent), and a separate errors list for
    values that failed (e.g., zero matching features).

    Chains into: ``qgis_figures_to_pptx``.
    """
    _stub("qgis_batch_render", "4 / Export & Batch & Delivery")


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_figures_to_pptx(
    figure_paths: Annotated[list[str], Field(description="Absolute paths to PNG/JPG figure files. One per slide.")],
    pptx_path: Annotated[str, Field(description="Absolute path for the output .pptx.")],
    layout: Annotated[Literal["title_and_image", "image_only", "two_column", "title_image_caption"], Field(description="Slide layout per figure.")] = "title_and_image",
    captions: Annotated[list[str] | None, Field(description="Optional per-slide captions. Must match length of figure_paths if given.")] = None,
    template_pptx: Annotated[str | None, Field(description="If given, slides are appended to this template; else a new blank deck is created.")] = None,
) -> PptxResult:
    """Drop figures into a PowerPoint deck. Delivery tool — closes the W17 loop.

    When to use: as the final step of a figure pipeline. After rendering
    one or more PNGs (via ``qgis_render_*`` or ``qgis_batch_render``), call
    this once to assemble them into a deck.

    Returns: ``PptxResult`` with pptx_path, slides added/total, and
    per-slide titles (None for slides without titles).
    """
    from pptx import Presentation
    from pptx.util import Inches

    # python-pptx default master layouts: 5 = Title Only, 6 = Blank.
    # title_image_caption / two_column degrade to title_only for v0.3.
    LAYOUT_INDEX = {
        "title_and_image": 5,
        "image_only": 6,
        "two_column": 5,
        "title_image_caption": 5,
    }

    if captions is not None and len(captions) != len(figure_paths):
        raise ValueError(
            f"captions length ({len(captions)}) must match figure_paths length "
            f"({len(figure_paths)}). Pass captions=None to skip titles entirely."
        )

    abs_pptx = os.path.abspath(pptx_path)
    prs = Presentation(os.path.abspath(template_pptx)) if template_pptx else Presentation()
    layout_idx = LAYOUT_INDEX.get(layout, 5)
    chosen_layout = prs.slide_layouts[layout_idx]

    slide_titles: list[str | None] = []
    for i, fig in enumerate(figure_paths):
        slide = prs.slides.add_slide(chosen_layout)
        title_text = None
        if captions is not None and layout_idx != 6 and slide.shapes.title is not None:
            title_text = captions[i]
            slide.shapes.title.text = title_text
        slide_titles.append(title_text)
        slide.shapes.add_picture(
            os.path.abspath(fig), Inches(0.5), Inches(1.5), height=Inches(5.5)
        )

    prs.save(abs_pptx)
    return PptxResult(
        pptx_path=abs_pptx,
        n_slides_added=len(figure_paths),
        n_slides_total=len(prs.slides),
        slide_titles=slide_titles,
    )


# ---------------------------------------------------------------------------
# Tools — Escape hatch (1)
# ---------------------------------------------------------------------------


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=False, destructiveHint=True, openWorldHint=True
    )
)
def qgis_eval(
    code: Annotated[str, Field(description="PyQGIS source to execute.")],
    return_vars: Annotated[list[str] | None, Field(description="Local variable names to capture from the executed scope and return JSON-serialized.")] = None,
) -> EvalResult:
    """Execute arbitrary PyQGIS. Escape hatch — prefer the workflow tools.

    When to use: only when no workflow tool fits — e.g., an unusual
    processing algorithm, a custom symbology that ``qgis_style_*`` can't
    express, or interactive debugging. For routine work (load, style,
    render, export, batch, deck) use the dedicated tools.

    In plugin transport, executes inside the running QGIS process. In
    headless transport, executes in the standalone PyQGIS subprocess.
    Either way, this can mutate state — annotated as destructive.

    Returns: ``EvalResult`` with stdout, stderr, captured ``return_values``
    (if ``return_vars`` was given), and exception traceback (if any).
    """
    _stub("qgis_eval", "4 / Escape hatch")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------


def main() -> None:
    """Run the MCP server over stdio (default)."""
    logger.info("qgis-mcp-north server starting (v0.3.0, 5 of 13 tools implemented)")
    mcp.run()


if __name__ == "__main__":
    main()
