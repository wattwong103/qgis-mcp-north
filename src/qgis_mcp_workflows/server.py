#!/usr/bin/env python3
"""qgis-mcp-workflows — focused QGIS MCP server for transportation research figures.

Forked from nkarasiak/qgis-mcp v0.2.1. Cut hard to 12 workflow tools + 1 escape
hatch (`qgis_eval`). See ``docs/DESIGN.md`` for tool-surface design, response
shapes, error taxonomy, and roadmap.

This module is the **v0.2 scaffold**: every tool is registered with FastMCP and
its full input/output schema, but every body raises ``NotImplementedError``.
Implementation lands tool-by-tool from v0.3 onward (see DESIGN.md §7).

Default socket port: 9877 (vs upstream nkarasiak/qgis-mcp on 9876). Both servers
can run side-by-side; the LLM picks per request based on tool descriptions.
"""

from __future__ import annotations

import logging
import os
import sys
from logging.handlers import RotatingFileHandler
from typing import Annotated, Literal

from mcp.server.fastmcp import FastMCP
from mcp.types import ToolAnnotations
from pydantic import BaseModel, Field

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------


def _setup_logging() -> logging.Logger:
    """stderr (WARNING+) plus optional rotating file handler."""
    log = logging.getLogger("QgisMcpWorkflowsServer")
    log.handlers.clear()

    fmt = logging.Formatter("%(asctime)s - %(name)s - %(levelname)s - %(message)s")

    stderr_handler = logging.StreamHandler(sys.stderr)
    stderr_handler.setLevel(logging.WARNING)
    stderr_handler.setFormatter(fmt)
    log.addHandler(stderr_handler)

    default_log_file = os.path.join("~", ".local", "share", "qgis-mcp-workflows", "server.log")
    log_file_raw = os.environ.get("QGIS_MCP_WORKFLOWS_LOG_FILE", default_log_file)
    log_level_name = os.environ.get("QGIS_MCP_WORKFLOWS_LOG_LEVEL", "INFO").upper()
    file_level = getattr(logging, log_level_name, logging.INFO)

    if log_file_raw:
        log_file = os.path.expanduser(log_file_raw)
        os.makedirs(os.path.dirname(log_file), exist_ok=True)
        file_handler = RotatingFileHandler(log_file, maxBytes=5 * 1024 * 1024, backupCount=3)
        file_handler.setLevel(file_level)
        file_handler.setFormatter(fmt)
        log.addHandler(file_handler)
        log.setLevel(min(logging.WARNING, file_level))
    else:
        log.setLevel(logging.WARNING)
    return log


logger = _setup_logging()


# ---------------------------------------------------------------------------
# Optional movingpandas — enables speed-binned trajectory rendering when
# installed via `uv sync --extra trajectory`. Detected once at module load;
# tests patch this attribute to exercise both code paths.
# ---------------------------------------------------------------------------

try:
    import movingpandas as _mp  # noqa: F401

    _HAS_MP = True
except Exception:  # ImportError or any movingpandas init failure
    _HAS_MP = False


# ---------------------------------------------------------------------------
# FastMCP server
# ---------------------------------------------------------------------------

SERVER_INSTRUCTIONS = """\
qgis-mcp-workflows — opinionated QGIS MCP for transportation research figure
pipelines (PFLOW, GUFM). Reach for these tools when you need to:
- Inspect or load a vector/raster layer.
- Render a choropleth from a zone polygon + value CSV.
- Render trajectories from CSV (PFLOW schema: lon, lat, datetime, trip_id).
- Render origin-destination flow arcs from an OD CSV + zones layer.
- Drop a batch of figures into a PowerPoint deck for a weekly report.
- Run arbitrary PyQGIS as an escape hatch (`qgis_eval`).

If you need general QGIS-API features (feature editing, processing
algorithms, layer-tree management, plugin tooling), use the upstream
nkarasiak/qgis-mcp server side-by-side. Both run together; pick per task.
"""

mcp = FastMCP("qgis-mcp-workflows", instructions=SERVER_INSTRUCTIONS)


# ---------------------------------------------------------------------------
# Tool registration mode — full (13 tools) vs compound (5 grouped tools).
#
# Read at module load. Tests patch this attribute to verify both surfaces.
# Compound mode collapses the surface to qgis_inspect / qgis_style / qgis_render /
# qgis_export / qgis_eval for token-constrained LLMs (Haiku, small open-weights).
# ---------------------------------------------------------------------------

TOOL_MODE = os.environ.get("QGIS_MCP_WORKFLOWS_TOOL_MODE", "full").lower()
if TOOL_MODE not in ("full", "compound"):
    logger.warning("Unknown QGIS_MCP_WORKFLOWS_TOOL_MODE=%r; defaulting to 'full'", TOOL_MODE)
    TOOL_MODE = "full"


def _maybe_tool(*args, **kwargs):
    """Conditional @mcp.tool decorator — registers when TOOL_MODE == 'full', else no-op.

    Functions decorated with this are still callable in Python; the only effect of
    no-op'ing is that FastMCP doesn't expose them over MCP. Direct imports (e.g.
    from tests) work the same in both modes.
    """
    if TOOL_MODE == "full":
        return mcp.tool(*args, **kwargs)
    return lambda f: f


def _maybe_compound_tool(*args, **kwargs):
    """Conditional decorator for compound-only tools — registers when TOOL_MODE == 'compound'."""
    if TOOL_MODE == "compound":
        return mcp.tool(*args, **kwargs)
    return lambda f: f


def _register_compound_tools_if_enabled() -> None:
    """Import compound.py at module-load tail to trigger _maybe_compound_tool decorators.

    Imports are conditional: only fire when TOOL_MODE='compound' to avoid pulling
    in compound.py's dependencies (which import every standalone tool) on full-mode
    cold-starts. Tests patch TOOL_MODE and reimport for both surfaces.
    """
    if TOOL_MODE == "compound":
        from qgis_mcp_workflows import compound  # noqa: F401  — import for side effects


# ---------------------------------------------------------------------------
# Pydantic response models — one per tool; see DESIGN.md §4 for field rationale
# ---------------------------------------------------------------------------


class FieldInfo(BaseModel):
    name: str
    type: str
    n_unique: int | None = None


class LayerInfo(BaseModel):
    """Read-only metadata about a vector or raster file on disk."""

    path: str
    geometry_type: Literal["point", "line", "polygon", "raster", "no_geom"]
    crs: str
    n_features: int
    extent: list[float] = Field(..., description="[xmin, ymin, xmax, ymax]")
    fields: list[FieldInfo]


class LoadedLayer(LayerInfo):
    """Layer that has been added to the active QGIS project."""

    layer_id: str


class LayerSummary(BaseModel):
    """Compact summary of a layer registered in a project."""

    layer_id: str
    name: str
    geometry_type: str
    visible: bool


class LayoutSummary(BaseModel):
    name: str


class ProjectInfo(BaseModel):
    project_path: str
    crs: str
    extent: list[float]
    layers: list[LayerSummary]
    layouts: list[LayoutSummary]


class ClassEntry(BaseModel):
    value: str
    color: str
    n_features: int


class StyleResult(BaseModel):
    layer_id: str
    n_classes: int
    classes: list[ClassEntry]


class GraduatedStyleResult(StyleResult):
    breaks: list[float]
    mode: str
    diverging: bool = False
    center: float = 0.0
    diverging_one_sided: bool = False


# ---------------------------------------------------------------------------
# Basemap tile presets — no-API-key XYZ providers drawn UNDER the data.
# Each entry: (url_template, attribution, zmax). Esri REST tiles use {z}/{y}/{x}
# order; that ordering is encoded in the template and passed to the plugin verbatim.
# URLs sourced from the xyzservices registry; copied (not imported) to keep the
# live-XYZ path dependency-free.
# ---------------------------------------------------------------------------

BasemapName = Literal["none", "positron", "dark_matter", "voyager", "osm", "esri_imagery"]

_BASEMAP_PRESETS: dict[str, tuple[str, str, int]] = {
    "positron": (
        "https://basemaps.cartocdn.com/light_all/{z}/{x}/{y}.png",
        "© OpenStreetMap contributors © CARTO",
        20,
    ),
    "dark_matter": (
        "https://basemaps.cartocdn.com/dark_all/{z}/{x}/{y}.png",
        "© OpenStreetMap contributors © CARTO",
        20,
    ),
    "voyager": (
        "https://basemaps.cartocdn.com/rastertiles/voyager/{z}/{x}/{y}.png",
        "© OpenStreetMap contributors © CARTO",
        20,
    ),
    "osm": (
        "https://tile.openstreetmap.org/{z}/{x}/{y}.png",
        "© OpenStreetMap contributors",
        19,
    ),
    "esri_imagery": (
        "https://server.arcgisonline.com/ArcGIS/rest/services/World_Imagery/MapServer/tile/{z}/{y}/{x}",
        "Esri, Maxar, Earthstar Geographics",
        19,
    ),
}


def _resolve_basemap(basemap: str, opacity: float) -> dict | None:
    """Resolve a basemap preset name into the ``basemap_spec`` sent to the plugin.

    ``"none"`` returns ``None`` (legacy white-background behavior, unchanged).
    A known preset returns a live-XYZ spec the plugin loads via
    ``QgsRasterLayer(type=xyz, provider="wms")``.
    """
    if basemap == "none":
        return None
    url, attribution, zmax = _BASEMAP_PRESETS[basemap]
    return {
        "kind": "xyz",
        "name": basemap,
        "url": url,
        "zmin": 0,
        "zmax": zmax,
        "attribution": attribution,
        "opacity": float(opacity),
    }


class RenderResult(BaseModel):
    output_path: str
    width: int
    height: int
    dpi: int
    extent: list[float]
    crs: str
    n_layers: int
    basemap_attribution: str | None = None
    basemap_source: str | None = None


class JoinResult(BaseModel):
    csv: str
    field: str
    n_matched: int
    n_unmatched: int


class ChoroplethResult(RenderResult):
    field: str
    n_classes: int
    breaks: list[float]
    mode: str
    min_value: float
    max_value: float
    n_features: int
    join: JoinResult | None = None
    diverging: bool = False
    center: float = 0.0
    diverging_one_sided: bool = False


class TrajectoryResult(RenderResult):
    n_trajectories: int
    n_points_total: int
    n_points_rendered: int
    downsampled: bool
    time_range: list[str] | None = None
    modes: list[str] | None = None
    used_movingpandas: bool


class ODFlowResult(RenderResult):
    n_flows: int
    n_flows_rendered: int
    n_zones: int
    max_flow: float
    min_flow_rendered: float
    n_unmatched_origins: int
    n_unmatched_destinations: int


class LinkDensityResult(RenderResult):
    n_trajectory_rows_total: int
    n_points_total: int
    n_links_with_traffic: int
    n_links_rendered: int
    n_unmatched_link_ids: int
    density_field: str
    breaks: list[float]
    mode: str
    min_density: float
    max_density: float
    aggregation: str


class ExportResult(BaseModel):
    output_path: str
    format: str
    n_pages: int
    layout_name: str


class ComposeLayoutResult(BaseModel):
    output_path: str
    format: str
    n_layers: int
    items: list[str]
    page_size_mm: list[float]


class DiagramMapResult(RenderResult):
    diagram_type: str
    value_fields: list[str]
    n_features: int


class CatchmentResult(RenderResult):
    method: str
    n_points: int
    n_catchments: int


class BatchManifestEntry(BaseModel):
    value: str
    output_path: str
    extent: list[float]


class BatchError(BaseModel):
    value: str
    error: str


class BatchRenderResult(BaseModel):
    output_dir: str
    n_rendered: int
    manifest: list[BatchManifestEntry]
    errors: list[BatchError]


class PptxResult(BaseModel):
    pptx_path: str
    n_slides_added: int
    n_slides_total: int
    slide_titles: list[str | None]


class EvalResult(BaseModel):
    stdout: str
    stderr: str
    return_values: dict | None = None
    exception: str | None = None


# ---------------------------------------------------------------------------
# Stub helper
# ---------------------------------------------------------------------------


def _stub(tool_name: str, design_section: str) -> None:
    """Raise a clear NotImplementedError pointing at the design doc."""
    raise NotImplementedError(
        f"{tool_name} is a v0.2 scaffold stub. Implementation lands in a future "
        f"version — see docs/DESIGN.md §{design_section} for the spec."
    )


_RASTER_EXTENSIONS = (".tif", ".tiff", ".geotiff", ".vrt", ".asc", ".img", ".jp2")


def _is_raster_path(path: str) -> bool:
    return path.lower().endswith(_RASTER_EXTENSIONS)


def _translate_geometry_type(plugin_type: str) -> str:
    """Plugin's ``vector_{0,1,2}`` / ``raster`` → DESIGN.md geometry_type enum."""
    if plugin_type == "raster":
        return "raster"
    if plugin_type.startswith("vector_"):
        idx = plugin_type.split("_", 1)[1]
        return {"0": "point", "1": "line", "2": "polygon"}.get(idx, "no_geom")
    return "no_geom"


def _load_and_get_info(executor, abs_path: str, name: str | None = None):
    """Load a layer + fetch metadata. Removes the layer if get_layer_info fails.

    Returns ``(layer_id, info_dict, is_raster)``. Layer stays loaded on success
    so the caller can decide whether to remove it (transient inspect) or keep it
    (persistent load).
    """
    is_raster = _is_raster_path(abs_path)
    load_cmd = "add_raster_layer" if is_raster else "add_vector_layer"
    params = {"path": abs_path}
    if name:
        params["name"] = name
    load_result = executor.dispatch(load_cmd, params)
    layer_id = load_result["id"]
    try:
        info = executor.dispatch("get_layer_info", {"layer_id": layer_id})
    except Exception:
        try:
            executor.dispatch("remove_layer", {"layer_id": layer_id})
        except Exception:
            logger.warning("post-error cleanup failed for %s", layer_id, exc_info=True)
        raise
    return layer_id, info, is_raster


def _layer_info_kwargs(abs_path: str, info: dict, is_raster: bool) -> dict:
    """Translate plugin's get_layer_info response into LayerInfo constructor kwargs."""
    extent_dict = info["extent"]
    fields = [
        FieldInfo(name=f["name"], type=f["type"], n_unique=None)
        for f in info.get("fields", [])
    ]
    return {
        "path": abs_path,
        "geometry_type": _translate_geometry_type(info["type"]),
        "crs": info["crs"],
        "n_features": 0 if is_raster else info.get("feature_count", 0),
        "extent": [
            extent_dict["xmin"], extent_dict["ymin"],
            extent_dict["xmax"], extent_dict["ymax"],
        ],
        "fields": fields,
    }


# ---------------------------------------------------------------------------
# Tools — Inspection & loading (3)
# ---------------------------------------------------------------------------


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=True, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_layer_inspect(
    path: Annotated[str, Field(description="Absolute path to a shapefile (.shp), GeoPackage (.gpkg), GeoJSON (.geojson), or raster (.tif).")],
) -> LayerInfo:
    """Read metadata for a layer file *without* loading it.

    When to use: before any render_* call, to confirm field names, geometry
    type, CRS, and feature count. Cheap (no project mutation, no QGIS render).

    Inputs: absolute path to a vector or raster file on disk.

    Returns: ``LayerInfo`` with path, geometry_type ∈ {point, line, polygon,
    raster, no_geom}, CRS string (EPSG:#### or proj string), n_features,
    extent as ``[xmin, ymin, xmax, ymax]``, and the list of fields with their
    types and unique-value counts.

    Chains into: ``qgis_load_layer`` (if you need to render),
    ``qgis_render_choropleth`` (pass ``zones_path=path``),
    ``qgis_render_trajectory`` (pass ``input_path=path``).
    """
    from qgis_mcp_workflows.executors import get_executor

    abs_path = os.path.abspath(path)
    executor = get_executor()
    layer_id, info, is_raster = _load_and_get_info(executor, abs_path)
    try:
        return LayerInfo(**_layer_info_kwargs(abs_path, info, is_raster))
    finally:
        try:
            executor.dispatch("remove_layer", {"layer_id": layer_id})
        except Exception:
            logger.warning("transient cleanup failed for %s", layer_id, exc_info=True)


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=False, destructiveHint=False, openWorldHint=True
    )
)
def qgis_load_layer(
    path: Annotated[str, Field(description="Absolute path to a vector or raster file.")],
    name: Annotated[str | None, Field(description="Display name for the layer in the project. Defaults to the file's stem.")] = None,
    crs: Annotated[str | None, Field(description='Override CRS, e.g. "EPSG:4326". Use when the file is missing a .prj sidecar.')] = None,
) -> LoadedLayer:
    """Add a layer to the active project and return its metadata + layer_id.

    When to use: when downstream tools need to operate on a registered layer
    (styling, rendering with multiple layers, batch operations). For one-shot
    figure renders on a single file, ``qgis_render_choropleth`` and friends
    accept a path directly and don't require this step.

    Returns: ``LoadedLayer`` — same fields as ``LayerInfo`` plus ``layer_id``
    used by ``qgis_style_*`` and ``qgis_render_map``.

    Chains into: ``qgis_style_categorized``, ``qgis_style_graduated``,
    ``qgis_render_map``.
    """
    from qgis_mcp_workflows.errors import CrsMismatchError, ExecutorError
    from qgis_mcp_workflows.executors import get_executor

    abs_path = os.path.abspath(path)
    executor = get_executor()
    layer_id, info, is_raster = _load_and_get_info(executor, abs_path, name=name)
    kwargs = _layer_info_kwargs(abs_path, info, is_raster)

    if crs is not None:
        try:
            crs_result = executor.dispatch("set_layer_crs", {"layer_id": layer_id, "crs": crs})
        except ExecutorError as err:
            try:
                executor.dispatch("remove_layer", {"layer_id": layer_id})
            except Exception:
                logger.warning("post-error cleanup failed for %s", layer_id, exc_info=True)
            raise CrsMismatchError(crs, err.message) from err
        kwargs["crs"] = crs_result.get("crs", crs)

    return LoadedLayer(layer_id=layer_id, **kwargs)


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=False, destructiveHint=False, openWorldHint=True
    )
)
def qgis_project_load(
    qgz_path: Annotated[str, Field(description="Absolute path to a .qgz or .qgs project file.")],
) -> ProjectInfo:
    """Load a saved QGIS project (``.qgz`` / ``.qgs``).

    When to use: when the user has already designed a map in QGIS — layers
    loaded, styles applied, layouts configured — and you want to export
    figures from it. The W17 weekly-deck pattern uses this: load project,
    then ``qgis_export_layout`` or ``qgis_render_map``.

    Returns: ``ProjectInfo`` listing all layers (with layer_ids) and all
    print-composer layouts available for export.

    Chains into: ``qgis_export_layout``, ``qgis_batch_render``,
    ``qgis_render_map``.
    """
    from qgis_mcp_workflows.errors import ExecutorError, ProjectLoadError
    from qgis_mcp_workflows.executors import get_executor

    abs_qgz = os.path.abspath(qgz_path)
    try:
        result = get_executor().dispatch("project_load", {"qgz_path": abs_qgz}, timeout=30)
    except ExecutorError as err:
        raise ProjectLoadError(abs_qgz, err.message) from err

    return ProjectInfo(
        project_path=result["project_path"],
        crs=result["crs"],
        extent=result["extent"],
        layers=[
            LayerSummary(
                layer_id=la["layer_id"],
                name=la["name"],
                geometry_type=la["geometry_type"],
                visible=la["visible"],
            )
            for la in result.get("layers", [])
        ],
        layouts=[LayoutSummary(name=lo["name"]) for lo in result.get("layouts", [])],
    )


# ---------------------------------------------------------------------------
# Tools — Styling (2)
# ---------------------------------------------------------------------------


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=False
    )
)
def qgis_style_categorized(
    layer_id: Annotated[str, Field(description="layer_id from qgis_load_layer or qgis_project_load.")],
    field: Annotated[str, Field(description="Field name to categorize on (string-typed).")],
    palette: Annotated[str, Field(description='ColorBrewer palette name, e.g. "Set2", "Paired", "Dark2".')] = "Set2",
    classes: Annotated[list[str] | None, Field(description="Optional subset/order of category values to render; others get a default 'no data' style.")] = None,
) -> StyleResult:
    """Apply categorical (one-color-per-value) symbology to a vector layer.

    When to use: when ``field`` holds a small set of categorical values
    (e.g., ``transport_mode``, ``taxi_type``, ``zone_type``) and you want
    each category in a distinct color.

    Returns: ``StyleResult`` with the resolved class list and per-class
    feature counts.
    """
    from qgis_mcp_workflows.errors import ExecutorError, FieldNotFoundError, LayerNotFoundError
    from qgis_mcp_workflows.executors import get_executor

    params = {
        "layer_id": layer_id,
        "style_type": "categorized",
        "field": field,
        "color_ramp": palette,
    }
    if classes is not None:
        params["classes_subset"] = list(classes)

    try:
        result = get_executor().dispatch("set_layer_style", params, timeout=30)
    except ExecutorError as err:
        if "Field not found" in err.message:
            raise FieldNotFoundError(field, []) from err
        if "Layer not found" in err.message:
            raise LayerNotFoundError(layer_id) from err
        raise

    return StyleResult(
        layer_id=layer_id,
        n_classes=result["n_classes"],
        classes=[
            ClassEntry(value=c["value"], color=c["color"], n_features=c["n_features"])
            for c in result.get("classes", [])
        ],
    )


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=False
    )
)
def qgis_style_graduated(
    layer_id: Annotated[str, Field(description="layer_id from qgis_load_layer or qgis_project_load.")],
    field: Annotated[str, Field(description="Numeric field to bin on (e.g., total_trips, trip_count, vkt).")],
    n_classes: Annotated[int, Field(description="Number of bins.", ge=2, le=15)] = 5,
    mode: Annotated[Literal["quantile", "equal_interval", "natural_breaks", "pretty"], Field(description="Binning strategy.")] = "quantile",
    palette: Annotated[str, Field(description='Sequential colorbrewer palette, e.g. "YlOrRd", "Blues", "Viridis".')] = "YlOrRd",
    diverging: Annotated[bool, Field(description="Diverging color scheme with a fixed neutral midpoint, for signed data (e.g. net flux). Replaces mode-based boundaries with symmetric breaks around center; pair with a diverging palette (vik/RdBu/balance).")] = False,
    center: Annotated[float, Field(description="Neutral midpoint for diverging mode (e.g. 0). Ignored when diverging is False.")] = 0.0,
) -> GraduatedStyleResult:
    """Apply graduated (value-based color ramp) symbology — the choropleth primitive.

    When to use: as a low-level building block. For zone-level choropleths,
    prefer ``qgis_render_choropleth`` (one call instead of three).

    Returns: ``GraduatedStyleResult`` — class list + per-class feature counts +
    the resolved ``breaks`` array + the ``mode`` used.
    """
    from qgis_mcp_workflows.errors import ExecutorError, FieldNotFoundError, LayerNotFoundError
    from qgis_mcp_workflows.executors import get_executor

    params = {
        "layer_id": layer_id,
        "style_type": "graduated",
        "field": field,
        "classes": n_classes,
        "mode": mode,
        "color_ramp": palette,
        "diverging": diverging,
        "center": center,
    }

    try:
        result = get_executor().dispatch("set_layer_style", params, timeout=30)
    except ExecutorError as err:
        if "Field not found" in err.message:
            raise FieldNotFoundError(field, []) from err
        if "Layer not found" in err.message:
            raise LayerNotFoundError(layer_id) from err
        raise

    return GraduatedStyleResult(
        layer_id=layer_id,
        n_classes=result["n_classes"],
        classes=[
            ClassEntry(value=c["value"], color=c["color"], n_features=c["n_features"])
            for c in result.get("classes", [])
        ],
        breaks=result.get("breaks", []),
        mode=result.get("mode", mode),
        diverging=result.get("diverging", diverging),
        center=result.get("center", center),
        diverging_one_sided=result.get("diverging_one_sided", False),
    )


# ---------------------------------------------------------------------------
# Tools — Rendering (4)
# ---------------------------------------------------------------------------


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_render_map(
    layer_ids: Annotated[list[str], Field(description="Layers to render, drawn in order bottom→top.")],
    output_png: Annotated[str, Field(description="Absolute path for the output PNG.")],
    width: Annotated[int, Field(description="Image width in pixels.", ge=200, le=8000)] = 1600,
    height: Annotated[int, Field(description="Image height in pixels.", ge=200, le=8000)] = 1200,
    dpi: Annotated[int, Field(description="Image DPI (affects font size).", ge=72, le=600)] = 150,
    extent: Annotated[list[float] | None, Field(description="Render extent as [xmin, ymin, xmax, ymax]. If omitted, uses the union of layer extents with 5% padding.")] = None,
    background: Annotated[str, Field(description='Map background, named or hex (e.g. "white", "#fafafa", "transparent").')] = "white",
) -> RenderResult:
    """Render a list of already-loaded layers to PNG.

    When to use: generic render. For domain-specific cases prefer the
    workflow tools (``qgis_render_choropleth``, ``qgis_render_trajectory``,
    ``qgis_render_od_flows``) — they handle data loading, styling, and
    extent inference for you.

    Returns: ``RenderResult`` with the absolute output_path, image dims,
    final extent, CRS, and number of layers rendered.

    Chains into: ``qgis_figures_to_pptx``, ``qgis_batch_render``.
    """
    from qgis_mcp_workflows.executors import get_executor

    abs_output = os.path.abspath(output_png)
    params: dict = {
        "layer_ids": list(layer_ids),
        "output_png": abs_output,
        "width": width,
        "height": height,
        "dpi": dpi,
        "background": background,
    }
    if extent is not None:
        params["extent"] = list(extent)

    result = get_executor().dispatch("render_layers_to_path", params, timeout=60)
    return RenderResult(
        output_path=result["output_path"],
        width=result["width"],
        height=result["height"],
        dpi=result["dpi"],
        extent=result["extent"],
        crs=result["crs"],
        n_layers=result["n_layers"],
    )


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_render_choropleth(
    zones_path: Annotated[str, Field(description="Absolute path to a polygon layer (shp, gpkg, geojson). For PFLOW prefecture choropleth, use polbnda_jpn_new.shp.")],
    value_field: Annotated[str, Field(description="Field name to render. If value_csv is given, this is the column in that CSV; else this is an attribute on zones_path.")],
    output_png: Annotated[str, Field(description="Absolute path for the output PNG.")],
    value_csv: Annotated[str | None, Field(description="Optional CSV path joined to zones_path on join_field. Use this for the PFLOW pattern: zone_trips.csv joined to a zones polygon.")] = None,
    join_field: Annotated[str, Field(description='Common key between zones_path and value_csv. PFLOW uses "zone_id" (MFS-coded), "PRF_CODE" (prefecture), or similar.')] = "zone_id",
    n_classes: Annotated[int, Field(description="Number of choropleth bins.", ge=2, le=15)] = 5,
    mode: Annotated[Literal["quantile", "equal_interval", "natural_breaks", "pretty"], Field(description="Binning strategy.")] = "quantile",
    palette: Annotated[str, Field(description='Sequential colorbrewer palette, e.g. "YlOrRd", "Blues", "Viridis".')] = "YlOrRd",
    diverging: Annotated[bool, Field(description="Diverging color scheme pinned at a neutral midpoint, for signed data (net flux = arrivals minus departures). Symmetric class breaks around center; pair with a diverging palette (vik/RdBu/balance).")] = False,
    center: Annotated[float, Field(description="Neutral midpoint for diverging mode (e.g. 0). Ignored when diverging is False.")] = 0.0,
    label_field: Annotated[str | None, Field(description="Optional zones attribute to label each polygon with (e.g. a ward/prefecture name), drawn with a white halo for legibility.")] = None,
    title: Annotated[str | None, Field(description="Optional title rendered at the top of the figure.")] = None,
    legend: Annotated[bool, Field(description="Render a legend with class breaks.")] = True,
    basemap_paths: Annotated[list[str] | None, Field(description="Optional vector basemap layers drawn under the choropleth (e.g., coastline, rivers, prefecture borders).")] = None,
    basemap: Annotated[BasemapName, Field(description='Tile basemap drawn under the data for real-world context. "positron"/"voyager" = neutral grey (best for choropleths), "dark_matter" = dark, "osm" = streets, "esri_imagery" = satellite. "none" keeps the legacy white background. No API key needed.')] = "none",
    basemap_opacity: Annotated[float, Field(description="Opacity of the tile basemap, 0.0–1.0. Use 0.5–0.8 to mute it so the choropleth colors read on top.", ge=0.0, le=1.0)] = 1.0,
    width: Annotated[int, Field(description="Image width in pixels.", ge=200, le=8000)] = 1600,
    height: Annotated[int, Field(description="Image height in pixels.", ge=200, le=8000)] = 1200,
    dpi: Annotated[int, Field(description="Image DPI.", ge=72, le=600)] = 150,
) -> ChoroplethResult:
    """Render a zone-level choropleth in one call. PFLOW workflow tool.

    When to use: any zone-aggregated visualization where each polygon gets
    a color from a numeric value. Replaces ``qgis_load_layer`` +
    ``qgis_style_graduated`` + ``qgis_render_map`` chain.

    Two data shapes supported:
    1. ``value_field`` is already an attribute on ``zones_path`` → render directly.
    2. ``value_csv`` is provided → left-join ``value_csv[join_field]`` to
       ``zones_path[join_field]``, then render. Mismatches surface in the
       response as ``join.n_unmatched``, never as silent zero values.

    PFLOW example: ``zones_path=polbnda_jpn_new.shp``,
    ``value_csv=zone_trips.csv``, ``value_field=total_trips``,
    ``join_field=zone_id``.

    Returns: ``ChoroplethResult`` with output path, breaks, min/max, and a
    ``join`` block (if value_csv was used) reporting matched vs unmatched.

    Chains into: ``qgis_figures_to_pptx``, ``qgis_batch_render``.
    """
    import csv as _csv

    from qgis_mcp_workflows.errors import ExecutorError, FieldNotFoundError, JoinError
    from qgis_mcp_workflows.executors import get_executor

    abs_zones = os.path.abspath(zones_path)
    abs_output = os.path.abspath(output_png)
    abs_basemaps = [os.path.abspath(p) for p in (basemap_paths or [])]
    abs_csv = os.path.abspath(value_csv) if value_csv else None

    value_dict: dict[str, float] | None = None
    if abs_csv is not None:
        with open(abs_csv, encoding="utf-8", newline="") as f:
            reader = _csv.DictReader(f)
            csv_columns = reader.fieldnames or []
            if value_field not in csv_columns:
                raise FieldNotFoundError(value_field, csv_columns)
            if join_field not in csv_columns:
                raise FieldNotFoundError(join_field, csv_columns)
            value_dict = {}
            for row in reader:
                key = row[join_field]
                raw = row[value_field]
                try:
                    value_dict[str(key)] = float(raw)
                except (TypeError, ValueError):
                    logger.warning(
                        "skipping non-numeric value in %s: %s=%r → %s=%r",
                        abs_csv, join_field, key, value_field, raw,
                    )

    params: dict = {
        "zones_path": abs_zones,
        "value_field": value_field,
        "output_png": abs_output,
        "value_dict": value_dict,
        "join_field": join_field,
        "n_classes": n_classes,
        "mode": mode,
        "palette": palette,
        "diverging": diverging,
        "center": center,
        "label_field": label_field,
        "basemap_paths": abs_basemaps,
        "basemap_spec": _resolve_basemap(basemap, basemap_opacity),
        "width": width,
        "height": height,
        "dpi": dpi,
    }

    try:
        result = get_executor().dispatch("render_choropleth", params, timeout=60)
    except ExecutorError as err:
        if "JOIN_NO_MATCH" in str(err):
            raise JoinError(err.message) from err
        raise

    join_block = None
    if value_dict is not None:
        join_block = JoinResult(
            csv=abs_csv or "",
            field=join_field,
            n_matched=result["n_matched"],
            n_unmatched=result["n_unmatched"],
        )
    return ChoroplethResult(
        output_path=result["output_path"],
        width=result["width"],
        height=result["height"],
        dpi=result["dpi"],
        extent=result["extent"],
        crs=result["crs"],
        n_layers=result["n_layers"],
        field=result["field"],
        n_classes=result["n_classes"],
        breaks=result["breaks"],
        mode=result["mode"],
        min_value=result["min_value"],
        max_value=result["max_value"],
        n_features=result["n_features"],
        join=join_block,
        diverging=result.get("diverging", diverging),
        center=result.get("center", center),
        diverging_one_sided=result.get("diverging_one_sided", False),
        basemap_attribution=result.get("basemap_attribution"),
        basemap_source=result.get("basemap_source"),
    )


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_render_trajectory(
    input_path: Annotated[str, Field(description="Absolute path to a CSV (with lon, lat, datetime, trip_id columns by default — matches PFLOW trajectory schema) or GPX file.")],
    output_png: Annotated[str, Field(description="Absolute path for the output PNG.")],
    lon_col: Annotated[str, Field(description="Longitude column name in the CSV.")] = "lon",
    lat_col: Annotated[str, Field(description="Latitude column name in the CSV.")] = "lat",
    time_col: Annotated[str, Field(description="Datetime column name (ISO 8601 or PFLOW 'YYYY-MM-DD HH:MM:SS' string).")] = "datetime",
    id_col: Annotated[str, Field(description="Trajectory grouping column. Each unique value is one trajectory.")] = "trip_id",
    mode_col: Annotated[str | None, Field(description='Optional categorical column to color trajectories by (e.g., "transport_mode").')] = None,
    render_mode: Annotated[Literal["lines", "points", "heatmap"], Field(description="Visualization style.")] = "lines",
    sample_rate: Annotated[float, Field(description="Fraction of points to keep (1.0 = all, 0.01 = every 100th).", gt=0.0, le=1.0)] = 1.0,
    max_points: Annotated[int, Field(description="Hard cap on rendered points; exceeded → automatic downsample with response flag.", ge=1000)] = 500_000,
    basemap_paths: Annotated[list[str] | None, Field(description="Optional basemap layers drawn under trajectories.")] = None,
    extent: Annotated[list[float] | None, Field(description="[lon_min, lat_min, lon_max, lat_max] in EPSG:4326. Clips before rendering.")] = None,
    width: Annotated[int, Field(description="Image width in pixels.", ge=200, le=8000)] = 1600,
    height: Annotated[int, Field(description="Image height in pixels.", ge=200, le=8000)] = 1200,
    dpi: Annotated[int, Field(description="Image DPI.", ge=72, le=600)] = 150,
) -> TrajectoryResult:
    """Render trajectory data from CSV/GPX. PFLOW/GUFM workflow tool.

    When to use: visualizing GPS-style point sequences. Defaults match
    PFLOW trajectory CSV schema (``lon, lat, datetime, trip_id,
    transport_mode``). PFLOW files are large (3M+ rows, ~1 GB each); use
    ``sample_rate`` and ``extent`` to keep renders fast.

    If ``movingpandas`` and the QGIS Trajectools plugin are installed,
    auto-uses them for richer rendering (speed bins, stop detection); else
    falls back to plain line/point rendering. Reports which path was taken
    via ``used_movingpandas``.

    Returns: ``TrajectoryResult`` with totals, rendered count, downsample
    flag, observed time range, modes seen.

    Chains into: ``qgis_figures_to_pptx``.
    """
    import csv as _csv

    from qgis_mcp_workflows.errors import EmptyAfterFilterError, FieldNotFoundError
    from qgis_mcp_workflows.executors import get_executor

    abs_input = os.path.abspath(input_path)
    abs_output = os.path.abspath(output_png)
    abs_basemaps = [os.path.abspath(p) for p in (basemap_paths or [])]

    # GPX path: skip CSV parse, hand path through to plugin's OGR loader.
    if abs_input.lower().endswith(".gpx"):
        params: dict = {
            "input_path": abs_input,
            "output_png": abs_output,
            "render_mode": render_mode,
            "basemap_paths": abs_basemaps,
            "extent": list(extent) if extent is not None else None,
            "width": width,
            "height": height,
            "dpi": dpi,
            "features": None,
            "mode_col": mode_col,
            "used_movingpandas": False,
            "speed_field": None,
        }
        result = get_executor().dispatch("render_trajectory", params, timeout=120)
        return TrajectoryResult(
            output_path=result["output_path"],
            width=result["width"], height=result["height"], dpi=result["dpi"],
            extent=result["extent"], crs=result["crs"], n_layers=result["n_layers"],
            n_trajectories=result["n_trajectories"],
            n_points_total=result["n_points_total"],
            n_points_rendered=result["n_points_rendered"],
            downsampled=result["downsampled"],
            time_range=result.get("time_range"),
            modes=result.get("modes"),
            used_movingpandas=result.get("used_movingpandas", False),
        )

    # CSV path: parse + validate columns MCP-side.
    with open(abs_input, encoding="utf-8", newline="") as f:
        reader = _csv.DictReader(f)
        columns = reader.fieldnames or []
        for required in (lon_col, lat_col, time_col, id_col):
            if required not in columns:
                raise FieldNotFoundError(required, columns)
        if mode_col is not None and mode_col not in columns:
            raise FieldNotFoundError(mode_col, columns)
        rows = list(reader)

    n_points_total = len(rows)

    # Apply extent clip first (caller's filter), then sampling.
    if extent is not None:
        xmin, ymin, xmax, ymax = extent
        kept = []
        for row in rows:
            try:
                lon = float(row[lon_col])
                lat = float(row[lat_col])
            except (TypeError, ValueError):
                continue
            if xmin <= lon <= xmax and ymin <= lat <= ymax:
                kept.append(row)
        rows = kept
        if not rows:
            raise EmptyAfterFilterError(
                f"0 rows after extent clip [{xmin}, {ymin}, {xmax}, {ymax}]"
            )

    # Sampling: stride by 1/sample_rate, then cap by max_points.
    downsampled = False
    if sample_rate < 1.0:
        stride = max(1, round(1.0 / sample_rate))
        rows = rows[::stride]
        if not rows:
            raise EmptyAfterFilterError(
                f"0 rows after sample_rate={sample_rate} (stride={stride})"
            )
    if len(rows) > max_points:
        stride2 = -(-len(rows) // max_points)  # ceil division
        rows = rows[::stride2]
        downsampled = True

    # Build feature list — small dicts, JSON-safe over the socket.
    features: list[dict] = []
    modes_seen: set[str] = set()
    time_min: str | None = None
    time_max: str | None = None
    for row in rows:
        try:
            lon = float(row[lon_col])
            lat = float(row[lat_col])
        except (TypeError, ValueError):
            continue
        ts = row.get(time_col, "")
        if time_min is None or (ts and ts < time_min):
            time_min = ts
        if time_max is None or (ts and ts > time_max):
            time_max = ts
        feat: dict = {
            "trip_id": str(row[id_col]),
            "lon": lon,
            "lat": lat,
            "datetime": ts,
        }
        if mode_col is not None:
            mode_val = row.get(mode_col, "")
            feat["mode"] = mode_val
            if mode_val:
                modes_seen.add(mode_val)
        features.append(feat)

    if not features:
        raise EmptyAfterFilterError("0 valid rows after numeric coercion")

    # movingpandas integration: speed-binned line rendering only.
    used_mp = False
    speed_field: str | None = None
    if _HAS_MP and render_mode == "lines" and mode_col is None:
        try:
            import movingpandas as mp  # uses sys.modules; tests patch this in
            speeds = _compute_movingpandas_speeds(mp, features)
            if speeds is not None and len(speeds) == len(features):
                for feat, sp in zip(features, speeds, strict=False):
                    feat["speed_kmh"] = sp
                used_mp = True
                speed_field = "speed_kmh"
        except Exception:
            logger.warning("movingpandas integration failed, falling back", exc_info=True)

    n_trajectories = len({f["trip_id"] for f in features})
    time_range = [time_min, time_max] if (time_min and time_max) else None
    modes_list = sorted(modes_seen) if modes_seen else None

    params = {
        "input_path": abs_input,
        "output_png": abs_output,
        "render_mode": render_mode,
        "basemap_paths": abs_basemaps,
        "extent": list(extent) if extent is not None else None,
        "width": width,
        "height": height,
        "dpi": dpi,
        "features": features,
        "mode_col": mode_col,
        "used_movingpandas": used_mp,
        "speed_field": speed_field,
    }
    result = get_executor().dispatch("render_trajectory", params, timeout=120)
    return TrajectoryResult(
        output_path=result["output_path"],
        width=result["width"], height=result["height"], dpi=result["dpi"],
        extent=result["extent"], crs=result["crs"], n_layers=result["n_layers"],
        n_trajectories=result.get("n_trajectories", n_trajectories),
        n_points_total=result.get("n_points_total", n_points_total),
        n_points_rendered=result.get("n_points_rendered", len(features)),
        downsampled=result.get("downsampled", downsampled),
        time_range=result.get("time_range", time_range),
        modes=result.get("modes", modes_list),
        used_movingpandas=result.get("used_movingpandas", used_mp),
    )


def _compute_movingpandas_speeds(mp, features: list[dict]) -> list[float] | None:
    """Build a TrajectoryCollection and return per-point speeds in km/h.

    Returns None if movingpandas can't produce a length-matched speed series.
    """
    try:
        import pandas as pd

        df = pd.DataFrame(features)
        df["datetime"] = pd.to_datetime(df["datetime"])
        try:
            import geopandas as gpd
            from shapely.geometry import Point

            gdf = gpd.GeoDataFrame(
                df,
                geometry=[Point(lon, lat) for lon, lat in zip(df["lon"], df["lat"], strict=False)],
                crs="EPSG:4326",
            )
            tc = mp.TrajectoryCollection(gdf, traj_id_col="trip_id", t="datetime")
        except Exception:
            # Tests patch in a fake mp.TrajectoryCollection that accepts anything;
            # real movingpandas requires geopandas, which is bundled with the
            # [trajectory] extra. Fall through to a permissive constructor for the
            # test path.
            tc = mp.TrajectoryCollection(df, traj_id_col="trip_id", t="datetime")
        try:
            tc.add_speed(overwrite=True, units=("km", "h"), name="speed_kmh")
        except TypeError:
            tc.add_speed()  # fakes may accept no kwargs
        point_gdf = tc.to_point_gdf()
        col = "speed_kmh" if "speed_kmh" in point_gdf.columns else "speed"
        if col not in point_gdf.columns:
            return None
        speeds = [float(v) if v == v else 0.0 for v in point_gdf[col].tolist()]
        return speeds
    except Exception:
        logger.warning("movingpandas speed computation failed", exc_info=True)
        return None


def _aggregate_link_density(
    csv_paths: list,
    link_id_col: str,
    aggregation: str,
    value_col: str | None,
) -> tuple[dict[str, float], int]:
    """Stream-aggregate trajectory CSVs into a {link_id → density} dict.

    Returns (density_dict, n_rows_read). Streaming: never holds more than one
    row in memory beyond the accumulator. Non-numeric values in value_col are
    skipped silently (logged at WARNING) so a few bad rows don't kill the run.

    Raises:
        FieldNotFoundError: if link_id_col or value_col is missing from any CSV.
        ValueError: if aggregation='sum' but value_col is None.
    """
    import csv as _csv

    from qgis_mcp_workflows.errors import FieldNotFoundError

    if aggregation == "sum" and value_col is None:
        raise ValueError("aggregation='sum' requires value_col to be set.")
    if aggregation not in ("count", "sum"):
        raise ValueError(f"Unknown aggregation: {aggregation!r}. Use 'count' or 'sum'.")

    density: dict[str, float] = {}
    n_rows = 0

    for path in csv_paths:
        with open(path, encoding="utf-8", newline="") as f:
            reader = _csv.DictReader(f)
            columns = reader.fieldnames or []
            if link_id_col not in columns:
                raise FieldNotFoundError(link_id_col, columns)
            if value_col is not None and value_col not in columns:
                raise FieldNotFoundError(value_col, columns)

            for row in reader:
                n_rows += 1
                link_id = row[link_id_col]
                if not link_id:
                    continue
                if aggregation == "count":
                    density[link_id] = density.get(link_id, 0.0) + 1.0
                else:  # sum
                    raw = row[value_col]  # type: ignore[index]
                    try:
                        val = float(raw)
                    except (TypeError, ValueError):
                        continue
                    if val != val:  # NaN check
                        continue
                    density[link_id] = density.get(link_id, 0.0) + val

    return density, n_rows


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_render_od_flows(
    od_csv: Annotated[str, Field(description="Absolute path to a long-format OD CSV. PFLOW schema: origin, destination, trip_count, avg_distance_km.")],
    zones_layer_path: Annotated[str, Field(description="Absolute path to a polygon layer with one feature per zone, keyed by zone_id_field.")],
    output_png: Annotated[str, Field(description="Absolute path for the output PNG.")],
    origin_col: Annotated[str, Field(description="Origin zone column in od_csv.")] = "origin",
    dest_col: Annotated[str, Field(description="Destination zone column in od_csv.")] = "destination",
    value_col: Annotated[str, Field(description="Flow magnitude column. Arc widths scale linearly with this.")] = "trip_count",
    zone_id_field: Annotated[str, Field(description="Zone identifier field on zones_layer_path. Must match origin_col / dest_col values.")] = "zone_id",
    top_n: Annotated[int | None, Field(description="Render only the top-N flows by value. None renders all matched flows.")] = None,
    arc_style: Annotated[Literal["line", "arrow", "curved"], Field(description='Arc rendering: "line" (straight, default), "arrow" (directional), or "curved" (directional bezier). Arrows/curves scale width + head with flow.')] = "line",
    basemap_paths: Annotated[list[str] | None, Field(description="Optional vector basemap layers drawn under arcs.")] = None,
    basemap: Annotated[BasemapName, Field(description='Tile basemap drawn under the arcs ("positron"/"voyager"/"dark_matter"/"osm"/"esri_imagery"). "none" keeps the legacy white background. No API key needed.')] = "none",
    basemap_opacity: Annotated[float, Field(description="Opacity of the tile basemap, 0.0-1.0.", ge=0.0, le=1.0)] = 1.0,
    width: Annotated[int, Field(description="Image width in pixels.", ge=200, le=8000)] = 1600,
    height: Annotated[int, Field(description="Image height in pixels.", ge=200, le=8000)] = 1200,
    dpi: Annotated[int, Field(description="Image DPI.", ge=72, le=600)] = 150,
) -> ODFlowResult:
    """Render origin-destination arcs over a zones layer. PFLOW workflow tool.

    When to use: any "flows between regions" visualization. PFLOW example:
    ``od_csv=truck/run_*/od_flows.csv``,
    ``zones_layer_path=polbnda_jpn_new.shp``, ``zone_id_field=PRF_CODE``.

    PFLOW reality: multiple zone-id systems coexist (MFS##, PRF##, Z##) in
    different files. The caller must align ``od_csv[origin_col / dest_col]``
    with ``zones_layer[zone_id_field]``. Mismatches surface as
    ``n_unmatched_origins`` / ``n_unmatched_destinations`` — loud, not
    silent zero-flow renders.

    Returns: ``ODFlowResult`` with rendered count, max flow, and unmatched
    counts on each side.

    Chains into: ``qgis_figures_to_pptx``.
    """
    import csv as _csv

    from qgis_mcp_workflows.errors import FieldNotFoundError
    from qgis_mcp_workflows.executors import get_executor

    abs_od = os.path.abspath(od_csv)
    abs_zones = os.path.abspath(zones_layer_path)
    abs_output = os.path.abspath(output_png)
    abs_basemaps = [os.path.abspath(p) for p in (basemap_paths or [])]
    basemap_spec = _resolve_basemap(basemap, basemap_opacity)

    with open(abs_od, encoding="utf-8", newline="") as f:
        reader = _csv.DictReader(f)
        columns = reader.fieldnames or []
        for required in (origin_col, dest_col, value_col):
            if required not in columns:
                raise FieldNotFoundError(required, columns)
        flows: list[dict] = []
        for row in reader:
            try:
                value = float(row[value_col])
            except (TypeError, ValueError):
                logger.warning("skipping non-numeric flow value: %s", row)
                continue
            flows.append({
                "origin": str(row[origin_col]),
                "destination": str(row[dest_col]),
                "value": value,
            })

    # Sort descending by value so plugin can compute max_flow up front and so
    # top_n picks the largest. Stable sort keeps tie-breaks deterministic.
    flows.sort(key=lambda f: f["value"], reverse=True)
    if top_n is not None and top_n > 0:
        flows = flows[:top_n]

    params = {
        "od_csv": abs_od,
        "zones_path": abs_zones,
        "output_png": abs_output,
        "flows": flows,
        "zone_id_field": zone_id_field,
        "arc_style": arc_style,
        "basemap_paths": abs_basemaps,
        "basemap_spec": basemap_spec,
        "width": width,
        "height": height,
        "dpi": dpi,
    }
    result = get_executor().dispatch("render_od_flows", params, timeout=60)
    return ODFlowResult(
        output_path=result["output_path"],
        width=result["width"], height=result["height"], dpi=result["dpi"],
        extent=result["extent"], crs=result["crs"], n_layers=result["n_layers"],
        n_flows=result["n_flows"],
        n_flows_rendered=result["n_flows_rendered"],
        n_zones=result["n_zones"],
        max_flow=result["max_flow"],
        min_flow_rendered=result["min_flow_rendered"],
        n_unmatched_origins=result["n_unmatched_origins"],
        n_unmatched_destinations=result["n_unmatched_destinations"],
        basemap_attribution=result.get("basemap_attribution"),
        basemap_source=result.get("basemap_source"),
    )


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_render_link_density(
    trajectory_csvs: Annotated[list[str], Field(description="One or more PFLOW trajectory CSV paths. Each must contain link_id_col. Streamed (not loaded fully); safe for multi-GB inputs.")],
    drm_network_path: Annotated[str, Field(description="Absolute path to the pre-built DRM network GeoPackage. Build once via scripts/build_drm_network.py.")],
    output_png: Annotated[str, Field(description="Absolute path for the output PNG.")],
    link_id_col: Annotated[str, Field(description="Trajectory CSV column joining to DRM link_id.")] = "link_id",
    aggregation: Annotated[Literal["count", "sum"], Field(description="Per-link aggregation. 'count' = number of trajectory points; 'sum' requires value_col.")] = "count",
    value_col: Annotated[str | None, Field(description="Numeric column to sum (only used when aggregation='sum'). Non-numeric / NaN values are skipped.")] = None,
    n_classes: Annotated[int, Field(description="Number of graduated bins for symbology.", ge=2, le=15)] = 7,
    mode: Annotated[Literal["quantile", "equal_interval", "natural_breaks", "pretty"], Field(description="Binning strategy for graduated styling.")] = "quantile",
    palette: Annotated[str, Field(description='Sequential colorbrewer palette, e.g. "YlOrRd", "Blues", "Viridis".')] = "YlOrRd",
    min_density: Annotated[float, Field(description="Drop links with density below this. Use to denoise rare-traffic links.", ge=0.0)] = 1.0,
    top_n: Annotated[int | None, Field(description="Render only the top-N densest links. None = all matched links.")] = None,
    extent: Annotated[list[float] | None, Field(description="Render extent [xmin, ymin, xmax, ymax] in EPSG:4326. If omitted, uses DRM layer extent.")] = None,
    basemap_paths: Annotated[list[str] | None, Field(description="Optional vector basemap layers drawn under links.")] = None,
    basemap: Annotated[BasemapName, Field(description='Tile basemap drawn under the links ("positron"/"voyager"/"dark_matter"/"osm"/"esri_imagery"). "none" keeps the legacy white background. No API key needed.')] = "none",
    basemap_opacity: Annotated[float, Field(description="Opacity of the tile basemap, 0.0-1.0.", ge=0.0, le=1.0)] = 1.0,
    width: Annotated[int, Field(description="Image width in pixels.", ge=200, le=8000)] = 1600,
    height: Annotated[int, Field(description="Image height in pixels.", ge=200, le=8000)] = 1200,
    dpi: Annotated[int, Field(description="Image DPI.", ge=72, le=600)] = 150,
) -> LinkDensityResult:
    """Render a DRM-link traffic-density choropleth from PFLOW trajectories. v2 workflow tool.

    When to use: visualizing where in the road network traffic concentrates,
    aggregated over potentially many GB of trajectory CSVs. Major upgrade
    over `qgis_render_trajectory`'s raw GPS scatter when you care about
    network-level density rather than individual paths.

    Prerequisite: `assets/drm_network.gpkg` must exist (one-time build via
    `scripts/build_drm_network.py`). The tool raises ``DRMNetworkNotFoundError``
    with the exact build command if missing.

    Big-data discipline: trajectory CSVs are streamed row-by-row, never
    loaded fully. Aggregation happens MCP-side; only the per-link totals
    (typically <100k entries) are sent to the plugin.

    Returns: ``LinkDensityResult`` — totals, matched/unmatched counts,
    resolved breaks, min/max density.

    Chains into: ``qgis_figures_to_pptx``.
    """
    from qgis_mcp_workflows.errors import (
        DRMNetworkNotFoundError,
        EmptyAfterFilterError,
    )
    from qgis_mcp_workflows.executors import get_executor

    abs_drm = os.path.abspath(drm_network_path)
    abs_output = os.path.abspath(output_png)
    abs_basemaps = [os.path.abspath(p) for p in (basemap_paths or [])]
    abs_csvs = [os.path.abspath(p) for p in trajectory_csvs]
    basemap_spec = _resolve_basemap(basemap, basemap_opacity)

    if not os.path.exists(abs_drm):
        raise DRMNetworkNotFoundError(abs_drm)

    density, n_rows_total = _aggregate_link_density(
        csv_paths=abs_csvs,
        link_id_col=link_id_col,
        aggregation=aggregation,
        value_col=value_col,
    )

    n_points_total = int(sum(density.values())) if aggregation == "count" else n_rows_total

    if min_density > 0.0:
        density = {k: v for k, v in density.items() if v >= min_density}

    if top_n is not None and top_n > 0:
        sorted_items = sorted(density.items(), key=lambda kv: kv[1], reverse=True)[:top_n]
        density = dict(sorted_items)

    if not density:
        raise EmptyAfterFilterError(
            f"0 links left after min_density={min_density}, top_n={top_n}. "
            f"Aggregated {len(density)} links from {n_rows_total} rows."
        )

    params: dict = {
        "density": density,
        "drm_network_path": abs_drm,
        "output_png": abs_output,
        "link_id_col": link_id_col,
        "aggregation": aggregation,
        "value_col": value_col,
        "n_classes": n_classes,
        "mode": mode,
        "palette": palette,
        "extent": list(extent) if extent is not None else None,
        "basemap_paths": abs_basemaps,
        "basemap_spec": basemap_spec,
        "width": width,
        "height": height,
        "dpi": dpi,
    }
    result = get_executor().dispatch("render_link_density", params, timeout=120)

    return LinkDensityResult(
        output_path=result["output_path"],
        width=result["width"], height=result["height"], dpi=result["dpi"],
        extent=result["extent"], crs=result["crs"], n_layers=result["n_layers"],
        n_trajectory_rows_total=n_rows_total,
        n_points_total=n_points_total,
        n_links_with_traffic=result["n_links_with_traffic"],
        n_links_rendered=result["n_links_rendered"],
        n_unmatched_link_ids=result["n_unmatched_link_ids"],
        density_field=result["density_field"],
        breaks=result["breaks"],
        mode=result["mode"],
        min_density=result["min_density"],
        max_density=result["max_density"],
        aggregation=aggregation,
        basemap_attribution=result.get("basemap_attribution"),
        basemap_source=result.get("basemap_source"),
    )


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_render_diagram_map(
    layer_path: Annotated[str, Field(description="Absolute path to a vector layer (polygons or points).")],
    value_fields: Annotated[list[str], Field(description="Numeric fields to chart per feature (>=1). Each becomes a pie slice / bar.")],
    output_png: Annotated[str, Field(description="Absolute path for the output PNG.")],
    diagram_type: Annotated[Literal["pie", "bar"], Field(description="Chart drawn on each feature.")] = "pie",
    size: Annotated[float, Field(description="Diagram size in mm.", gt=0.0, le=80.0)] = 10.0,
    palette: Annotated[str, Field(description='Palette for the series colors (e.g. "Set2", "Dark2", "viridis").')] = "Set2",
    extent: Annotated[list[float] | None, Field(description="Render extent [xmin, ymin, xmax, ymax] in the layer's CRS. Omit for full extent + 5%.")] = None,
    basemap: Annotated[BasemapName, Field(description='Tile basemap under the diagrams. "none" = white background.')] = "none",
    basemap_opacity: Annotated[float, Field(description="Tile basemap opacity 0.0-1.0.", ge=0.0, le=1.0)] = 1.0,
    width: Annotated[int, Field(description="Image width in pixels.", ge=200, le=8000)] = 1600,
    height: Annotated[int, Field(description="Image height in pixels.", ge=200, le=8000)] = 1200,
    dpi: Annotated[int, Field(description="Image DPI.", ge=72, le=600)] = 150,
) -> DiagramMapResult:
    """Render pie/bar charts on each map feature — "chart in map". v2 workflow tool.

    When to use: show a small multivariate breakdown (e.g., arrivals vs departures,
    mode split) per zone/station directly on the map, instead of a single
    choropleth color. Each value_field becomes a pie slice or bar.

    Returns: ``DiagramMapResult`` with the diagram type, charted fields, and
    feature count.

    Chains into: ``qgis_compose_layout``, ``qgis_figures_to_pptx``.
    """
    from qgis_mcp_workflows.executors import get_executor

    abs_layer = os.path.abspath(layer_path)
    abs_output = os.path.abspath(output_png)
    params = {
        "layer_path": abs_layer,
        "value_fields": list(value_fields),
        "output_png": abs_output,
        "diagram_type": diagram_type,
        "size": size,
        "palette": palette,
        "extent": list(extent) if extent is not None else None,
        "basemap_spec": _resolve_basemap(basemap, basemap_opacity),
        "width": width,
        "height": height,
        "dpi": dpi,
    }
    result = get_executor().dispatch("render_diagram_map", params, timeout=120)
    return DiagramMapResult(
        output_path=result["output_path"],
        width=result["width"], height=result["height"], dpi=result["dpi"],
        extent=result["extent"], crs=result["crs"], n_layers=result["n_layers"],
        diagram_type=result["diagram_type"],
        value_fields=result["value_fields"],
        n_features=result["n_features"],
        basemap_attribution=result.get("basemap_attribution"),
        basemap_source=result.get("basemap_source"),
    )


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_render_catchment(
    points_path: Annotated[str, Field(description="Absolute path to a point layer (e.g. stations).")],
    output_png: Annotated[str, Field(description="Absolute path for the output PNG.")],
    method: Annotated[Literal["voronoi"], Field(description="Catchment method. 'voronoi' = Thiessen service areas (one cell per point, nearest-point tessellation). Buffer rings / network isochrones are future methods.")] = "voronoi",
    extent: Annotated[list[float] | None, Field(description="Render extent [xmin, ymin, xmax, ymax] in the layer's CRS. Omit for full extent + 5%.")] = None,
    basemap: Annotated[BasemapName, Field(description='Tile basemap under the catchments. "none" = white background.')] = "none",
    basemap_opacity: Annotated[float, Field(description="Tile basemap opacity 0.0-1.0.", ge=0.0, le=1.0)] = 1.0,
    width: Annotated[int, Field(description="Image width in pixels.", ge=200, le=8000)] = 1600,
    height: Annotated[int, Field(description="Image height in pixels.", ge=200, le=8000)] = 1200,
    dpi: Annotated[int, Field(description="Image DPI.", ge=72, le=600)] = 150,
) -> CatchmentResult:
    """Render Voronoi service-area catchments around points. v2 workflow tool.

    When to use: approximate each station/facility's service area as its Thiessen
    cell (nearest-point tessellation) — the catchment method from the TransInfor
    fig05. Uses a pure QgsGeometry Voronoi op (no Processing framework).

    Returns: ``CatchmentResult`` with the method, point count, and catchment count.

    Chains into: ``qgis_compose_layout``, ``qgis_figures_to_pptx``.
    """
    from qgis_mcp_workflows.executors import get_executor

    abs_points = os.path.abspath(points_path)
    abs_output = os.path.abspath(output_png)
    params = {
        "points_path": abs_points,
        "output_png": abs_output,
        "method": method,
        "extent": list(extent) if extent is not None else None,
        "basemap_spec": _resolve_basemap(basemap, basemap_opacity),
        "width": width,
        "height": height,
        "dpi": dpi,
    }
    result = get_executor().dispatch("render_catchment", params, timeout=180)
    return CatchmentResult(
        output_path=result["output_path"],
        width=result["width"], height=result["height"], dpi=result["dpi"],
        extent=result["extent"], crs=result["crs"], n_layers=result["n_layers"],
        method=result["method"],
        n_points=result["n_points"],
        n_catchments=result["n_catchments"],
        basemap_attribution=result.get("basemap_attribution"),
        basemap_source=result.get("basemap_source"),
    )


# ---------------------------------------------------------------------------
# Tools — Export & Batch & Delivery (3)
# ---------------------------------------------------------------------------


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_export_layout(
    qgz_path: Annotated[str, Field(description="Absolute path to a .qgz or .qgs project containing the layout.")],
    layout_name: Annotated[str, Field(description="Print-composer layout name. List available layouts via qgis_project_load.")],
    output_path: Annotated[str, Field(description="Absolute path for the output file. Extension can be .png, .pdf, or .svg.")],
    format: Annotated[Literal["png", "pdf", "svg"], Field(description="Output format. Should match output_path extension.")] = "png",
    dpi: Annotated[int, Field(description="Export DPI.", ge=72, le=600)] = 300,
) -> ExportResult:
    """Export a saved print-composer layout to file.

    When to use: when the user has designed a polished figure layout in
    QGIS (titles, legends, scale bars, multi-panel) and wants to export it
    rather than re-render programmatically.

    Returns: ``ExportResult`` with format, page count, and resolved layout name.

    Chains into: ``qgis_figures_to_pptx``, ``qgis_batch_render``.
    """
    from qgis_mcp_workflows.errors import ExecutorError, LayoutNotFoundError
    from qgis_mcp_workflows.executors import get_executor

    abs_qgz = os.path.abspath(qgz_path)
    abs_output = os.path.abspath(output_path)
    params = {
        "qgz_path": abs_qgz,
        "layout_name": layout_name,
        "output_path": abs_output,
        "format": format,
        "dpi": dpi,
    }
    try:
        result = get_executor().dispatch("export_layout", params, timeout=60)
    except ExecutorError as err:
        if "LAYOUT_NOT_FOUND" in err.message:
            import re

            avail_match = re.search(r"Available:\s*\[([^\]]*)\]", err.message)
            avail = []
            if avail_match:
                avail = [
                    s.strip().strip("'\"")
                    for s in avail_match.group(1).split(",")
                    if s.strip()
                ]
            raise LayoutNotFoundError(layout_name, avail) from err
        raise

    return ExportResult(
        output_path=result["output_path"],
        format=result["format"],
        n_pages=result["n_pages"],
        layout_name=result["layout_name"],
    )


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_compose_layout(
    layer_paths: Annotated[list[str], Field(description="Absolute paths to layers in bottom->top draw order (vector or raster). Styling comes from each layer's own saved style; default symbology otherwise.")],
    output_path: Annotated[str, Field(description="Absolute output path; format inferred from extension (.png / .pdf / .svg).")],
    title: Annotated[str | None, Field(description="Optional title across the top of the page.")] = None,
    extent: Annotated[list[float] | None, Field(description="Map extent [xmin, ymin, xmax, ymax] in the layers' CRS. If omitted, uses the union of layer extents + 5% padding.")] = None,
    page: Annotated[Literal["a4_landscape", "a4_portrait", "a3_landscape", "square"], Field(description="Page size preset.")] = "a4_landscape",
    legend: Annotated[bool, Field(description="Add a legend linked to the map.")] = True,
    scale_bar: Annotated[bool, Field(description="Add a scale bar linked to the map.")] = True,
    north_arrow: Annotated[bool, Field(description="Add a north arrow from QGIS's bundled SVGs.")] = True,
    dpi: Annotated[int, Field(description="Export DPI.", ge=72, le=600)] = 300,
) -> ComposeLayoutResult:
    """Compose a deck-ready print layout from layers and export it. v2 workflow tool.

    When to use: turn one or more data/rendered layers into a publication figure
    with a titled map panel plus a linked legend, scale bar and north arrow — the
    gap that ``qgis_export_layout`` (which only exports pre-authored .qgz layouts)
    leaves open. Single panel for now; multi-panel / inset is a future extension.

    Returns: ``ComposeLayoutResult`` with output path, format, layer count, the
    furniture items added, and page size in mm.

    Chains into: ``qgis_figures_to_pptx``.
    """
    from qgis_mcp_workflows.executors import get_executor

    abs_layers = [os.path.abspath(p) for p in layer_paths]
    abs_output = os.path.abspath(output_path)
    params = {
        "layer_paths": abs_layers,
        "output_path": abs_output,
        "title": title,
        "extent": list(extent) if extent is not None else None,
        "page": page,
        "legend": legend,
        "scale_bar": scale_bar,
        "north_arrow": north_arrow,
        "dpi": dpi,
    }
    result = get_executor().dispatch("compose_layout", params, timeout=120)
    return ComposeLayoutResult(
        output_path=result["output_path"],
        format=result["format"],
        n_layers=result["n_layers"],
        items=result["items"],
        page_size_mm=result["page_size_mm"],
    )


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_batch_render(
    template_qgz: Annotated[str, Field(description="Absolute path to a template project. Must contain a single 'active' layer to filter and (optionally) a layout to export.")],
    attribute: Annotated[str, Field(description="Field on the active layer to filter by.")],
    values: Annotated[list[str], Field(description="Filter values to iterate. One render per value.")],
    output_dir: Annotated[str, Field(description="Absolute path to a directory where renders are written.")],
    layout_name: Annotated[str | None, Field(description="If given, exports the layout with this name; otherwise renders the map canvas.")] = None,
    filename_template: Annotated[str, Field(description="Filename template, e.g. '{value}.png' or 'choropleth_{value}.png'.")] = "{value}.png",
) -> BatchRenderResult:
    """Fan-out: render the same template per filter value. Workflow tool.

    When to use: 'render the OD map for each scenario', 'render the
    choropleth for each timestep', 'one figure per prefecture'. Emits a
    manifest you can feed straight into ``qgis_figures_to_pptx``.

    Returns: ``BatchRenderResult`` with output_dir, n_rendered, manifest
    of (value → output_path → extent), and a separate errors list for
    values that failed (e.g., zero matching features).

    Chains into: ``qgis_figures_to_pptx``.
    """
    from qgis_mcp_workflows.errors import ExecutorError, FieldNotFoundError
    from qgis_mcp_workflows.executors import get_executor

    abs_template = os.path.abspath(template_qgz)
    abs_output_dir = os.path.abspath(output_dir)

    if not values:
        return BatchRenderResult(
            output_dir=abs_output_dir,
            n_rendered=0,
            manifest=[],
            errors=[],
        )

    params = {
        "template_qgz": abs_template,
        "attribute": attribute,
        "values": list(values),
        "output_dir": abs_output_dir,
        "layout_name": layout_name,
        "filename_template": filename_template,
    }
    try:
        result = get_executor().dispatch("batch_render", params, timeout=300)
    except ExecutorError as err:
        if "FIELD_NOT_FOUND" in err.message:
            import re

            field_match = re.search(r"FIELD_NOT_FOUND:\s*['\"]([^'\"]+)['\"]", err.message)
            avail_match = re.search(r"Available:\s*\[([^\]]*)\]", err.message)
            field = field_match.group(1) if field_match else attribute
            avail = []
            if avail_match:
                avail = [
                    s.strip().strip("'\"")
                    for s in avail_match.group(1).split(",")
                    if s.strip()
                ]
            raise FieldNotFoundError(field, avail) from err
        raise

    return BatchRenderResult(
        output_dir=result["output_dir"],
        n_rendered=result["n_rendered"],
        manifest=[
            BatchManifestEntry(
                value=m["value"], output_path=m["output_path"], extent=m["extent"]
            )
            for m in result.get("manifest", [])
        ],
        errors=[
            BatchError(value=e["value"], error=e["error"])
            for e in result.get("errors", [])
        ],
    )


@_maybe_tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=True, destructiveHint=False, openWorldHint=True
    )
)
def qgis_figures_to_pptx(
    figure_paths: Annotated[list[str], Field(description="Absolute paths to PNG/JPG figure files. One per slide.")],
    pptx_path: Annotated[str, Field(description="Absolute path for the output .pptx.")],
    layout: Annotated[Literal["title_and_image", "image_only", "two_column", "title_image_caption"], Field(description="Slide layout per figure.")] = "title_and_image",
    captions: Annotated[list[str] | None, Field(description="Optional per-slide captions. Must match length of figure_paths if given.")] = None,
    template_pptx: Annotated[str | None, Field(description="If given, slides are appended to this template; else a new blank deck is created.")] = None,
) -> PptxResult:
    """Drop figures into a PowerPoint deck. Delivery tool — closes the W17 loop.

    When to use: as the final step of a figure pipeline. After rendering
    one or more PNGs (via ``qgis_render_*`` or ``qgis_batch_render``), call
    this once to assemble them into a deck.

    Returns: ``PptxResult`` with pptx_path, slides added/total, and
    per-slide titles (None for slides without titles).
    """
    from pptx import Presentation
    from pptx.util import Inches

    # python-pptx default master layouts: 5 = Title Only, 6 = Blank.
    # title_image_caption / two_column degrade to title_only for v0.3.
    LAYOUT_INDEX = {
        "title_and_image": 5,
        "image_only": 6,
        "two_column": 5,
        "title_image_caption": 5,
    }

    if captions is not None and len(captions) != len(figure_paths):
        raise ValueError(
            f"captions length ({len(captions)}) must match figure_paths length "
            f"({len(figure_paths)}). Pass captions=None to skip titles entirely."
        )

    abs_pptx = os.path.abspath(pptx_path)
    prs = Presentation(os.path.abspath(template_pptx)) if template_pptx else Presentation()
    layout_idx = LAYOUT_INDEX.get(layout, 5)
    chosen_layout = prs.slide_layouts[layout_idx]

    slide_titles: list[str | None] = []
    for i, fig in enumerate(figure_paths):
        slide = prs.slides.add_slide(chosen_layout)
        title_text = None
        if captions is not None and layout_idx != 6 and slide.shapes.title is not None:
            title_text = captions[i]
            slide.shapes.title.text = title_text
        slide_titles.append(title_text)
        slide.shapes.add_picture(
            os.path.abspath(fig), Inches(0.5), Inches(1.5), height=Inches(5.5)
        )

    prs.save(abs_pptx)
    return PptxResult(
        pptx_path=abs_pptx,
        n_slides_added=len(figure_paths),
        n_slides_total=len(prs.slides),
        slide_titles=slide_titles,
    )


# ---------------------------------------------------------------------------
# Tools — Escape hatch (1)
# ---------------------------------------------------------------------------


@mcp.tool(
    annotations=ToolAnnotations(
        readOnlyHint=False, idempotentHint=False, destructiveHint=True, openWorldHint=True
    )
)
def qgis_eval(
    code: Annotated[str, Field(description="PyQGIS source to execute.")],
    return_vars: Annotated[list[str] | None, Field(description="Local variable names to capture from the executed scope and return JSON-serialized.")] = None,
) -> EvalResult:
    """Execute arbitrary PyQGIS. Escape hatch — prefer the workflow tools.

    When to use: only when no workflow tool fits — e.g., an unusual
    processing algorithm, a custom symbology that ``qgis_style_*`` can't
    express, or interactive debugging. For routine work (load, style,
    render, export, batch, deck) use the dedicated tools.

    In plugin transport, executes inside the running QGIS process. In
    headless transport, executes in the standalone PyQGIS subprocess.
    Either way, this can mutate state — annotated as destructive.

    Returns: ``EvalResult`` with stdout, stderr, captured ``return_values``
    (if ``return_vars`` was given), and exception traceback (if any).
    """
    from qgis_mcp_workflows.executors import get_executor

    params: dict = {"code": code}
    if return_vars is not None:
        params["return_vars"] = list(return_vars)

    result = get_executor().dispatch("execute_code", params, timeout=300)

    exception_text: str | None = None
    if not result.get("executed", True):
        exception_text = result.get("traceback") or result.get("error")

    return EvalResult(
        stdout=result.get("stdout", ""),
        stderr=result.get("stderr", ""),
        return_values=result.get("return_values") if return_vars is not None else None,
        exception=exception_text,
    )


# Trigger compound-mode tool registration if env var requested it.
# This must come AFTER all standalone tool functions are defined (compound.py imports them).
_register_compound_tools_if_enabled()


# ---------------------------------------------------------------------------
# Transport selection
# ---------------------------------------------------------------------------


def _plugin_reachable(host: str, port: int, timeout_s: float = 0.5) -> bool:
    """Probe ``host:port`` with a short connect timeout. Used by ``auto`` mode."""
    import socket as _socket

    try:
        with _socket.create_connection((host, port), timeout=timeout_s):
            return True
    except OSError:
        return False


def _build_executor(transport: str):
    """Construct the executor for the chosen transport.

    ``auto`` first probes the plugin port; falls back to headless if the plugin
    is not reachable. Errors from headless construction propagate so the user
    sees a single clear ``HeadlessUnavailableError`` rather than a silent fall.
    """
    from qgis_mcp_workflows.executors.plugin import PluginExecutor
    from qgis_mcp_workflows.helpers import DEFAULT_HOST, DEFAULT_PORT

    host = os.environ.get("QGIS_MCP_WORKFLOWS_HOST", DEFAULT_HOST)
    port = int(os.environ.get("QGIS_MCP_WORKFLOWS_PORT", str(DEFAULT_PORT)))

    if transport == "plugin":
        return PluginExecutor(host=host, port=port), "plugin"
    if transport == "headless":
        from qgis_mcp_workflows.executors.headless import HeadlessExecutor

        return HeadlessExecutor(), "headless"
    if transport == "auto":
        if _plugin_reachable(host, port):
            return PluginExecutor(host=host, port=port), "plugin"
        from qgis_mcp_workflows.executors.headless import HeadlessExecutor

        return HeadlessExecutor(), "headless"
    raise ValueError(f"Unknown transport: {transport!r}. Use plugin / headless / auto.")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------


def main() -> None:
    """Run the MCP server. CLI flag ``--transport`` overrides the env default."""
    import argparse

    parser = argparse.ArgumentParser(prog="qgis-mcp-workflows-server")
    parser.add_argument(
        "--transport",
        choices=("plugin", "headless", "auto"),
        default=os.environ.get("QGIS_MCP_WORKFLOWS_TRANSPORT", "auto"),
        help="QGIS backend: plugin (TCP socket to running QGIS Desktop), "
        "headless (PyQGIS subprocess), auto (probe plugin port, fall back to "
        "headless). Default: auto. Env: QGIS_MCP_WORKFLOWS_TRANSPORT.",
    )
    args = parser.parse_args()

    from qgis_mcp_workflows.executors import set_executor

    executor, chosen = _build_executor(args.transport)
    set_executor(executor)
    logger.info("qgis-mcp-workflows server starting (v1.3.0, transport=%s)", chosen)
    mcp.run()


if __name__ == "__main__":
    main()
