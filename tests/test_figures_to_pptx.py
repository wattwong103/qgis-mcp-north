"""Strict-TDD tests for qgis_figures_to_pptx — pure python-pptx, no plugin."""

from __future__ import annotations

import os
from pathlib import Path

import pytest


def _make_png(path: Path, color=(180, 60, 60), size=(120, 90)) -> Path:
    """Pillow ships transitively with python-pptx — generate real PNGs cheaply."""
    from PIL import Image

    Image.new("RGB", size, color=color).save(path)
    return path


# T1 — minimal create
def test_creates_pptx_with_one_slide_per_image(tmp_path: Path):
    from qgis_mcp_north.server import qgis_figures_to_pptx

    img = _make_png(tmp_path / "fig1.png")
    out = tmp_path / "out.pptx"

    result = qgis_figures_to_pptx(figure_paths=[str(img)], pptx_path=str(out))

    assert out.exists(), "pptx file should be saved at the requested path"
    assert result.pptx_path == os.path.abspath(str(out))
    assert result.n_slides_added == 1
    assert result.n_slides_total == 1


# T2 — title_and_image layout uses captions[i] as slide title
def test_title_and_image_uses_caption_as_title(tmp_path: Path):
    from pptx import Presentation

    from qgis_mcp_north.server import qgis_figures_to_pptx

    img = _make_png(tmp_path / "fig.png")
    out = tmp_path / "out.pptx"

    result = qgis_figures_to_pptx(
        figure_paths=[str(img)],
        pptx_path=str(out),
        layout="title_and_image",
        captions=["Prefecture Total Trips"],
    )

    assert result.slide_titles == ["Prefecture Total Trips"]
    prs = Presentation(str(out))
    title_shape = prs.slides[0].shapes.title
    assert title_shape is not None, "title_and_image layout must have a title placeholder"
    assert title_shape.text == "Prefecture Total Trips"


# T3 — captions length must match figure_paths length
def test_captions_wrong_length_raises_value_error(tmp_path: Path):
    from qgis_mcp_north.server import qgis_figures_to_pptx

    img1 = _make_png(tmp_path / "f1.png")
    img2 = _make_png(tmp_path / "f2.png")
    out = tmp_path / "out.pptx"

    with pytest.raises(ValueError, match="captions"):
        qgis_figures_to_pptx(
            figure_paths=[str(img1), str(img2)],
            pptx_path=str(out),
            captions=["only_one"],  # mismatch
        )


# T4 — template_pptx appends new slides to an existing deck
def test_template_pptx_appends_slides(tmp_path: Path):
    from pptx import Presentation

    from qgis_mcp_north.server import qgis_figures_to_pptx

    template = tmp_path / "tpl.pptx"
    seed = Presentation()
    seed.slides.add_slide(seed.slide_layouts[5])  # 1 pre-existing slide
    seed.save(str(template))

    img = _make_png(tmp_path / "fig.png")
    out = tmp_path / "out.pptx"
    result = qgis_figures_to_pptx(
        figure_paths=[str(img)],
        pptx_path=str(out),
        template_pptx=str(template),
    )

    assert result.n_slides_added == 1
    assert result.n_slides_total == 2, "1 from template + 1 added"
    prs = Presentation(str(out))
    assert len(prs.slides) == 2


# T5 — missing figure_path raises actionable FileNotFoundError before pptx logic
def test_missing_figure_path_raises_file_not_found(tmp_path: Path):
    from qgis_mcp_north.server import qgis_figures_to_pptx

    out = tmp_path / "out.pptx"
    with pytest.raises(FileNotFoundError, match="does_not_exist.png"):
        qgis_figures_to_pptx(
            figure_paths=[str(tmp_path / "does_not_exist.png")],
            pptx_path=str(out),
        )
    assert not out.exists(), "must not produce a partial pptx when inputs are bad"
